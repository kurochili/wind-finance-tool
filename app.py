"""
Wind Farm Financial Assessment Dashboard
风电项目经济性评估看板 — 多项目管理 + 分项编辑 + 对比

启动方式: streamlit run app.py
"""

import copy
import io
import json
import os
import time
import uuid
from dataclasses import asdict
from typing import Dict, Literal, Optional, Tuple

import numpy as np
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
import streamlit as st

from wind_finance.calculator import CalculationResult, calculate
from wind_finance.country_profiles import (
    CountryProfile,
    CountryOMDefaults,
    MarketBenchmark,
    get_country_profile,
    get_diag_thresholds,
    list_countries,
    _PROFILES,
)
from wind_finance.excel_export import export_to_excel
try:
    from wind_finance.smart_input import (
        parse_excel, parse_image, extract_mw_from_name,
        detect_country_from_excel, detect_country_from_image_text,
    )
    _HAS_SMART_INPUT = True
except ImportError:
    _HAS_SMART_INPUT = False
from wind_finance.models import (
    BOPCost,
    BasicInfo,
    FinancingTerms,
    FoundationCost,
    InstallationCost,
    InvestmentData,
    OEMCost,
    OffshoreEPCBreakdown,
    OffshoreExtraCost,
    OnshoreInvestment,
    OperationalCost,
    PostWarrantyPeriodCost,
    TaxAndFinancial,
    WarrantyPeriodCost,
    WindFarmFinancialInputs,
    OM_METHODS,
    OM_METHOD_LABELS,
    OM_METHOD_DESCRIPTIONS,
)
from wind_finance.reverse_solver import (
    solve_hours_for_target_irr,
    solve_investment_for_target_lcoe,
    solve_tariff_for_target_irr,
    solve_tariff_for_zero_npv,
    solve_turbine_price_for_target_lcoe,
    solve_turbine_price_for_target_irr,
)
from wind_finance import db as _db
from wind_finance.tariff_data import get_tariff_references, get_tariff_display, get_all_tariff_summary

# ════════════════════════════════════════════════════════════════════════════
# 页面配置 & 全局样式
# ════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Wind Farm Assessment",
    page_icon="🌬️",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── 密码保护 ──────────────────────────────────────────────────────────────
# 密码来源优先级: HF Secrets > 环境变量 APP_PASSWORD > 不设密码（本地开发）
def _get_secret(key, default=""):
    val = os.environ.get(key, "")
    if not val:
        try:
            val = st.secrets.get(key, default)
        except Exception:
            val = default
    return val

_APP_PASSWORD = _get_secret("APP_PASSWORD")

if _APP_PASSWORD:
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False

    if not st.session_state.authenticated:
        st.markdown(
            "<div style='max-width:400px;margin:120px auto;text-align:center;'>"
            "<h2 style='color:#1F4E79;'>Wind Farm Financial Assessment</h2>"
            "<p style='color:#666;'>Please enter the access password.</p></div>",
            unsafe_allow_html=True,
        )
        col_l, col_m, col_r = st.columns([1, 1, 1])
        with col_m:
            pwd_input = st.text_input("Password", type="password", label_visibility="collapsed",
                                      placeholder="Enter password...")
            if st.button("Login", type="primary", use_container_width=True):
                if pwd_input == _APP_PASSWORD:
                    st.session_state.authenticated = True
                    st.rerun()
                else:
                    st.error("Wrong password. Please try again.")
        st.stop()

# ── 全局样式 ──────────────────────────────────────────────────────────────
st.markdown("""
<style>
    div[data-testid="stSidebar"] { background: #f8f9fb; }
    .block-container { padding-top: 1.5rem; }
    h1 { color: #1F4E79; }
    h2, h3 { color: #2E75B6; }
</style>
""", unsafe_allow_html=True)

COLOR_PALETTE = [
    "#1F4E79", "#2E75B6", "#548235", "#BF8F00",
    "#C00000", "#7030A0", "#ED7D31", "#4472C4",
    "#70AD47", "#FFC000", "#5B9BD5", "#A5A5A5",
]

# ════════════════════════════════════════════════════════════════════════════
# Supabase 初始化
# ════════════════════════════════════════════════════════════════════════════

_sb_url = _get_secret("SUPABASE_URL",
                       "https://vqahmhvxnjrxduwbxfzj.supabase.co")
_sb_key = _get_secret("SUPABASE_KEY",
                       "sb_publishable_dwrsDYm2eMjkZZP-hq9D_Q__OfPfHe3")
if _sb_url and _sb_key:
    _db.init(_sb_url, _sb_key)
_USE_DB = _db.db_available()

# ════════════════════════════════════════════════════════════════════════════
# Session State: 多项目存储（优先从数据库加载）
# ════════════════════════════════════════════════════════════════════════════

def _preload_defaults() -> Dict[str, dict]:
    """从内置脚本加载预置项目。"""
    import importlib, logging
    _log = logging.getLogger("preload")
    projects: Dict[str, dict] = {}
    _preloaders = [
        "wind_finance.preload_philippines",
        "wind_finance.preload_laguna",
        "wind_finance.preload_fsg",
        "wind_finance.preload_vietnam_qh",
        "wind_finance.preload_soctrang",
        "wind_finance.preload_hatinh",
    ]
    for _mod_name in _preloaders:
        try:
            _mod = importlib.import_module(_mod_name)
            entries = _mod.get_all_projects()
        except Exception as e:
            _log.warning("Failed to import %s: %s", _mod_name, e)
            continue
        for entry in entries:
            try:
                name, group, country, inputs, result = entry
                pid = str(uuid.uuid4())[:8]
                projects[pid] = {
                    "name": name,
                    "group": group,
                    "country": country,
                    "inputs": copy.deepcopy(inputs),
                    "result": copy.deepcopy(result),
                    "saved_at": time.strftime("%Y-%m-%d %H:%M:%S"),
                }
            except Exception as e:
                _log.warning("Failed to load project from %s: %s", _mod_name, e)
    _log.info("Preloaded %d projects from %d modules", len(projects), len(_preloaders))
    return projects


if "projects" not in st.session_state:
    st.session_state.projects: Dict[str, dict] = {}
    if _USE_DB:
        try:
            db_projects = _db.db_load_all()
            if db_projects:
                st.session_state.projects = db_projects
            else:
                defaults = _preload_defaults()
                st.session_state.projects = defaults
                for pid, proj in defaults.items():
                    try:
                        _db.db_save(pid, proj["name"], proj["group"],
                                    proj["country"], proj["inputs"], proj["saved_at"])
                    except Exception:
                        pass
        except Exception:
            st.session_state.projects = _preload_defaults()
    else:
        st.session_state.projects = _preload_defaults()

if "compare_ids" not in st.session_state:
    st.session_state.compare_ids: list[str] = []


def save_project(name: str, inputs: WindFarmFinancialInputs, result: CalculationResult,
                  group: str = "", country: str = "") -> str:
    pid = str(uuid.uuid4())[:8]
    _group = group or inputs.basic.project_name.split(" - ")[0].strip()
    _country = country or inputs.basic.country
    saved_at = time.strftime("%Y-%m-%d %H:%M:%S")
    st.session_state.projects[pid] = {
        "name": name,
        "group": _group,
        "country": _country,
        "inputs": copy.deepcopy(inputs),
        "result": copy.deepcopy(result),
        "saved_at": saved_at,
    }
    if _USE_DB:
        try:
            _db.db_save(pid, name, _group, _country, inputs, saved_at)
        except Exception:
            pass
    return pid


_DELETE_USER = _get_secret("DELETE_USER", "admin")
_DELETE_PWD  = _get_secret("DELETE_PWD")


def delete_project(pid: str):
    st.session_state.projects.pop(pid, None)
    if pid in st.session_state.compare_ids:
        st.session_state.compare_ids.remove(pid)
    st.session_state.pop("confirm_delete", None)
    if _USE_DB:
        try:
            _db.db_delete(pid)
        except Exception:
            pass


# ════════════════════════════════════════════════════════════════════════════
# 智能上传面板
# ════════════════════════════════════════════════════════════════════════════

def _generate_template_excel() -> bytes:
    """生成 Excel 参考模板"""
    buf = io.BytesIO()
    data = {
        "WTG Type": ["MySE8.5-230", "MySE10-242", "MySE16.X-260"],
        "Units": [47, 40, 25],
        "P90 [h/y]": [2070, 1944, 3200],
        "发电量提升比例": ["-", "-5.95%", "+8.2%"],
        "Electricity Tariff [USD/kWh]": [0.085, 0.085, 0.085],
        "RNA [USD/kW]": [355, 397, 520],
        "Tower [USD/kW]": [92, 53, 110],
        "Transportation [USD/kW]": [63, 49, 75],
        "Installation [USD/kW]": [101, 93, 130],
        "TSI [USD/kW]": [611, 592, 835],
        "BOP [USD/kW]": [839, 811, 650],
        "CAPEX [USD/kW]": [1450, 1403, 1485],
        "LCOE [USD/kWh]": ["", "", ""],
    }
    df = pd.DataFrame(data).T
    df.columns = [f"方案{i+1}" for i in range(df.shape[1])]
    df.index.name = "参数"
    df.to_excel(buf, engine="openpyxl")
    return buf.getvalue()


def _safe_float(val) -> Optional[float]:
    """安全转换为 float，无法转换则返回 None"""
    if val is None:
        return None
    try:
        f = float(val)
        return f if f != 0 else None
    except (ValueError, TypeError):
        return None


def _validate_smart_params(df: pd.DataFrame, tariff_default: float, is_offshore: bool):
    """
    校验智能输入表格，返回 (errors, warnings)。
    errors: 必须修正才能计算
    warnings: 风险提示，确认后可继续
    """
    errors = []
    warnings = []

    for idx, row in df.iterrows():
        label = row["方案"] if pd.notna(row["方案"]) else f"方案{idx+1}"

        # ── 错误：必填字段缺失或无效 ──
        units = row.get("台数")
        if pd.isna(units) or units is None or units <= 0:
            errors.append(f"**{label}** — 「台数」缺失或为 0，必须填写")
        mw = row.get("单机MW")
        if pd.isna(mw) or mw is None or mw <= 0:
            errors.append(f"**{label}** — 「单机MW」缺失或为 0，必须填写")
        p90 = row.get("P90(h)")
        if pd.isna(p90) or p90 is None or p90 <= 0:
            errors.append(f"**{label}** — 「P90等效小时数」缺失或为 0，必须填写")

        capex = row.get("CAPEX($/kW)")
        tsi = row.get("TSI($/kW)")
        bop = row.get("BOP($/kW)")
        has_capex = pd.notna(capex) and capex is not None and capex > 0
        has_tsi = pd.notna(tsi) and tsi is not None and tsi > 0
        has_bop = pd.notna(bop) and bop is not None and bop > 0
        if not has_capex and not (has_tsi and has_bop):
            errors.append(f"**{label}** — 投资数据缺失：请填写「CAPEX」或同时填写「TSI + BOP」")

        # ── 风险：数值异常范围 ──
        if pd.notna(p90) and p90 is not None and p90 > 0:
            if p90 < 1500:
                warnings.append(f"**{label}** — P90={int(p90)}h 偏低（通常 > 2000h），请确认")
            elif p90 > 5000:
                warnings.append(f"**{label}** — P90={int(p90)}h 偏高（通常 < 4500h），请确认")

        eff_capex = capex if has_capex else ((tsi or 0) + (bop or 0))
        if eff_capex and eff_capex > 0:
            if is_offshore and eff_capex < 500:
                warnings.append(f"**{label}** — CAPEX={eff_capex:.0f} $/kW 对海上项目偏低（通常 > 800），请确认")
            elif not is_offshore and eff_capex < 300:
                warnings.append(f"**{label}** — CAPEX={eff_capex:.0f} $/kW 偏低（通常 > 500），请确认")
            if eff_capex > 4000:
                warnings.append(f"**{label}** — CAPEX={eff_capex:.0f} $/kW 偏高（通常 < 3000），请确认")

        if pd.notna(mw) and mw is not None and mw > 0:
            if mw > 20:
                warnings.append(f"**{label}** — 单机 {mw:.1f}MW 超大容量，请确认是否正确")

        if pd.notna(units) and units is not None and units > 0:
            if units > 200:
                warnings.append(f"**{label}** — 台数 {int(units)} 台偏多，请确认")

        tariff_val = row.get("电价(USD/kWh)")
        if pd.notna(tariff_val) and tariff_val is not None and tariff_val > 0:
            if tariff_val > 0.25:
                warnings.append(f"**{label}** — 电价 {tariff_val:.4f} USD/kWh 偏高（> 0.25），请确认")
            elif tariff_val < 0.03:
                warnings.append(f"**{label}** — 电价 {tariff_val:.4f} USD/kWh 偏低（< 0.03），请确认")

        if has_tsi and has_bop and has_capex:
            if abs((tsi + bop) - capex) > 50:
                warnings.append(f"**{label}** — TSI({tsi:.0f}) + BOP({bop:.0f}) = {tsi+bop:.0f} 与 CAPEX({capex:.0f}) 差异较大，将以 CAPEX 为准")

    return errors, warnings


# ════════════════════════════════════════════════════════════════════════════
# 公共 O&M 参数渲染（三模式共用）
# ════════════════════════════════════════════════════════════════════════════

def _render_om_params(
    key_prefix: str,
    is_offshore: bool,
    profile,
    operation_years: int,
    total_investment_usd: float = 0.0,
    country_display: str = "",
    compact: bool = False,
    use_sidebar: bool = False,
) -> dict:
    """Render O&M parameter UI shared across Quick/Detailed/Smart Upload modes.

    Returns dict with keys: om_method, warranty, post_warranty,
    base_om_per_kw, om_escalation_rate, capex_om_percentage,
    capex_om_escalation, contract_om_periods.
    """
    _st = st.sidebar if use_sidebar else st
    om_d = getattr(profile, 'om_defaults', None) or CountryOMDefaults()
    _rec = om_d.recommended_method
    _opts_rev = {v: k for k, v in OM_METHOD_LABELS.items()}
    _rec_label = OM_METHOD_LABELS.get(_rec, "Fixed Escalation")

    if use_sidebar:
        _st.markdown("---")
        _st.markdown("### 🔄 运维计算方法")

    om_display = _st.selectbox(
        "O&M Algorithm",
        list(OM_METHOD_LABELS.values()),
        index=list(OM_METHOD_LABELS.keys()).index(_rec) if _rec in OM_METHOD_LABELS else 1,
        key=f"{key_prefix}_om_sel",
    )
    om_method = _opts_rev[om_display]
    _st.caption(f"🌏 {country_display} 推荐: **{_rec_label}**")
    _st.caption(OM_METHOD_DESCRIPTIONS.get(om_method, ""))

    _base_def = om_d.offshore_base_om if is_offshore else om_d.onshore_base_om
    _esc_def = om_d.escalation_rate
    _cpx_def = om_d.offshore_capex_pct if is_offshore else om_d.onshore_capex_pct

    warranty_cost = WarrantyPeriodCost()
    post_warranty_cost = PostWarrantyPeriodCost()
    base_om = float(_base_def)
    esc_rate = float(_esc_def)
    cpx_pct = float(_cpx_def)
    cpx_esc = 0.0
    ct_periods = [(1, 5, 15.0), (6, 10, 22.0), (11, operation_years, 28.0)]

    _k = key_prefix

    if om_method == "chinese_feasibility":
        _exp = _st.expander("🛡️ 质保期内运维 (USD/kW·年)", expanded=not compact)
        with _exp:
            wy = st.number_input("质保期 (年)", 0, 10, 5, key=f"{_k}_wy")
            wm = st.number_input("材料费", 0.0, 20.0, 4.23 if is_offshore else 2.82, step=0.1, key=f"{_k}_wm")
            wr = st.number_input("维修费", 0.0, 20.0, 0.0, step=0.1, key=f"{_k}_wr")
            wo = st.number_input("其他费用", 0.0, 20.0, 4.23 if is_offshore else 2.82, step=0.1, key=f"{_k}_wo")
            st.caption("⚪ 质保期成本由厂商承担大部分，对IRR影响较小。")
        warranty_cost = WarrantyPeriodCost(
            warranty_years=wy, material_cost_per_kw=wm,
            repair_cost_per_kw=wr, other_cost_per_kw=wo,
        )

        _exp2 = _st.expander("🔧 质保期外运维 (投资%递增)", expanded=not compact)
        with _exp2:
            pw_major = st.checkbox("含大部件更换", value=True, key=f"{_k}_pw_major")
            pw_mat = st.number_input("材料费 (USD/kW·年)", 0.0, 20.0,
                                      4.23 if is_offshore else 2.82, step=0.1, key=f"{_k}_pwm")
            pw_oth = st.number_input("其他费用 (USD/kW·年)", 0.0, 20.0,
                                      4.23 if is_offshore else 2.82, step=0.1, key=f"{_k}_pwo")

            st.markdown("**逐年维护费率（可逐行编辑）**")
            st.caption("日常维护 + 大部件替换均以静态投资为基数（%/年）")

            _stage_defs = [(1,5,0.5),(6,10,1.0),(11,15,1.5),(16,20,2.0),(21,25,2.5)]
            _tbl_key = f"{_k}_om_yr_table"
            _yrs_key = f"{_k}_om_yr_table_op_years"
            if _tbl_key not in st.session_state or st.session_state.get(_yrs_key) != operation_years:
                _rows = []
                for _yr in range(1, operation_years + 1):
                    _daily = 0.5
                    for _s, _e, _r in _stage_defs:
                        if _s <= _yr <= _e:
                            _daily = _r
                            break
                    _rows.append({"年份": _yr, "日常维护(%)": _daily, "大部件替换(%)": 0.0})
                st.session_state[_tbl_key] = pd.DataFrame(_rows)
                st.session_state[_yrs_key] = operation_years

            _om_edited = st.data_editor(
                st.session_state[_tbl_key],
                key=f"{_k}_om_yr_editor",
                use_container_width=True,
                hide_index=True,
                column_config={
                    "年份": st.column_config.NumberColumn("年", disabled=True, width="small"),
                    "日常维护(%)": st.column_config.NumberColumn("日常(%)", min_value=0.0, max_value=10.0, step=0.1, format="%.2f", width="small"),
                    "大部件替换(%)": st.column_config.NumberColumn("大部件(%)", min_value=0.0, max_value=10.0, step=0.1, format="%.2f", width="small"),
                },
                num_rows="fixed",
            )
            st.session_state[_tbl_key] = _om_edited

            if total_investment_usd > 0:
                _preview = []
                for _, _r in _om_edited.iterrows():
                    _tr = (_r["日常维护(%)"] + _r["大部件替换(%)"]) / 100.0
                    _ac = total_investment_usd * _tr / 10000
                    _preview.append(f"Y{int(_r['年份'])}: {_r['日常维护(%)']:.2f}%+{_r['大部件替换(%)']:.2f}% = {_ac:,.1f} 万USD")
                with st.expander("📊 逐年费用预览", expanded=False):
                    for _line in _preview:
                        st.caption(_line)

            maint_rates = []
            for _, _r in _om_edited.iterrows():
                _yr = int(_r["年份"])
                _tr = (_r["日常维护(%)"] + _r["大部件替换(%)"]) / 100.0
                maint_rates.append((_yr, _yr, _tr))

        post_warranty_cost = PostWarrantyPeriodCost(
            includes_major_components=pw_major,
            material_cost_per_kw=pw_mat,
            other_cost_per_kw=pw_oth,
            maintenance_rates=maint_rates,
        )

    elif om_method == "fixed_escalation":
        _exp = _st.expander("📊 固定单价法参数", expanded=not compact)
        with _exp:
            base_om = st.number_input("基准O&M (USD/kW/年)", 5.0, 100.0,
                                       float(_base_def), step=1.0, key=f"{_k}_fe_base")
            esc_rate = st.number_input("年增长率 (%)", 0.0, 10.0,
                                        float(_esc_def * 100), step=0.5, key=f"{_k}_fe_esc") / 100.0

    elif om_method == "capex_percentage":
        _exp = _st.expander("📊 投资百分比法参数", expanded=not compact)
        with _exp:
            cpx_pct = st.number_input("年运维费 (占CAPEX %)", 0.5, 10.0,
                                       float(_cpx_def * 100), step=0.1, key=f"{_k}_cp_pct") / 100.0
            cpx_esc = st.number_input("年增长率 (%)", 0.0, 10.0, 0.0,
                                       step=0.5, key=f"{_k}_cp_esc") / 100.0

    elif om_method == "contract":
        _exp = _st.expander("📋 合同报价法 (分阶段$/kW)", expanded=not compact)
        with _exp:
            st.caption("自定义各阶段运维单价")
            c1, c2, c3 = st.columns(3)
            cs1 = c1.number_input("阶段1起", 1, 35, 1, key=f"{_k}_cs1")
            ce1 = c2.number_input("阶段1止", 1, 35, 5, key=f"{_k}_ce1")
            cv1 = c3.number_input("$/kW/年", 0.0, 200.0, 15.0, step=1.0, key=f"{_k}_cv1")
            c1, c2, c3 = st.columns(3)
            cs2 = c1.number_input("阶段2起", 1, 35, 6, key=f"{_k}_cs2")
            ce2 = c2.number_input("阶段2止", 1, 35, 10, key=f"{_k}_ce2")
            cv2 = c3.number_input("$/kW/年", 0.0, 200.0, 22.0, step=1.0, key=f"{_k}_cv2")
            c1, c2, c3 = st.columns(3)
            cs3 = c1.number_input("阶段3起", 1, 35, 11, key=f"{_k}_cs3")
            ce3 = c2.number_input("阶段3止", 1, 35, operation_years, key=f"{_k}_ce3")
            cv3 = c3.number_input("$/kW/年", 0.0, 200.0, 28.0, step=1.0, key=f"{_k}_cv3")
            ct_periods = [(cs1, ce1, cv1), (cs2, ce2, cv2), (cs3, ce3, cv3)]

    return {
        "om_method": om_method,
        "warranty": warranty_cost,
        "post_warranty": post_warranty_cost,
        "base_om_per_kw": base_om,
        "om_escalation_rate": esc_rate,
        "capex_om_percentage": cpx_pct,
        "capex_om_escalation": cpx_esc,
        "contract_om_periods": ct_periods,
    }


def smart_upload_panel():
    """上传 Excel 或图片，自动提取参数并批量计算多方案"""
    st.markdown("### 📤 智能上传")
    st.caption("支持多方案自动识别 — 上传含 2~10 个机型的表格，系统自动提取每列为一个方案")

    if not _HAS_SMART_INPUT:
        st.error("智能输入模块尚未加载，请刷新页面重试。")
        return None

    upload_type = st.radio("文件类型", ["Excel (.xlsx)", "图片截图"], horizontal=True, key="su_type")

    if upload_type.startswith("Excel"):
        with st.expander("📋 Excel 格式说明 & 下载模板", expanded=False):
            st.markdown("""
**不需要严格模板！** 系统会自动识别以下关键字段行：

| 必填字段 | 可选字段 |
|---------|---------|
| WTG Type (机型名) | RNA, Tower, Transportation |
| Units (台数) | Installation, Tariff |
| P90 (等效小时数) | LCOE |
| TSI, BOP, CAPEX | 发电量提升比例 |

**要求**：
- 第一列是参数名，后续每列是一个方案（支持 2~10 个方案）
- 表格中可有空行、标题行，系统会自动跳过
- 数值支持逗号分隔(1,450)和百分号(-5.95%)
            """)
            tpl_bytes = _generate_template_excel()
            st.download_button(
                "📥 下载参考模板 Excel",
                data=tpl_bytes,
                file_name="wind_farm_input_template.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key="su_tpl_dl",
            )

        uploaded = st.file_uploader("上传 Excel 文件", type=["xlsx", "xls"], key="su_excel")
        if uploaded:
            raw_bytes = uploaded.read()
            variants = parse_excel(raw_bytes)
            _detected = detect_country_from_excel(raw_bytes, uploaded.name) if _HAS_SMART_INPUT else None
            if _detected:
                st.session_state["_su_detected_country"] = _detected
        else:
            variants = []
    else:
        st.info("支持截图自动识别：系统通过 OCR 读取表格中的数值，按列自动区分多个方案。建议截图清晰、表格规整。")
        uploaded = st.file_uploader("上传截图", type=["png", "jpg", "jpeg"], key="su_img")
        if uploaded:
            st.image(uploaded, caption="上传的截图", use_container_width=True)
            uploaded.seek(0)
            raw_bytes = uploaded.read()
            result = parse_image(raw_bytes)
            variants, _ocr_text, _ocr_err = result[0], result[1], result[2]
            if _ocr_err:
                st.error(f"图片识别失败：{_ocr_err}")
                variants = []
            elif not variants:
                st.warning("OCR 识别了文字但未能解析出参数。")
                with st.expander("🔍 查看 OCR 原始识别结果（调试用）", expanded=True):
                    st.code(_ocr_text, language=None)
                    st.caption("如果上方文字中能看到参数，说明是解析逻辑问题，请改用 Excel 上传或反馈给开发者。")
            else:
                with st.expander("🔍 查看 OCR 原始识别结果", expanded=False):
                    st.code(_ocr_text, language=None)
            _detected = detect_country_from_image_text(_ocr_text, uploaded.name) if _HAS_SMART_INPUT else None
            if _detected:
                st.session_state["_su_detected_country"] = _detected
        else:
            variants = []

    if not variants and not uploaded:
        return None

    if variants:
        st.success(f"成功提取 **{len(variants)}** 个方案：" +
                   " vs ".join(v.get("wtg_type", f"方案{i+1}") for i, v in enumerate(variants)))
    elif uploaded:
        st.info("未能自动提取方案，请在下方表格中**手动填写**参数。")
        variants = [{"wtg_type": "方案1"}, {"wtg_type": "方案2"}]

    # 国家自动检测
    countries = list_countries()
    country_options = {f"{cn} ({en})": en for en, cn in countries}
    detected_country = st.session_state.get("_su_detected_country")
    default_idx = 0
    if detected_country:
        for i, (display, en) in enumerate(country_options.items()):
            if en.lower() == detected_country.lower():
                default_idx = i
                break

    col_c, col_t = st.columns(2)
    with col_c:
        sel_country = st.selectbox("国家", list(country_options.keys()), index=default_idx, key="su_country")
        country_name = country_options[sel_country]
        if detected_country:
            st.caption(f"🔍 系统自动识别为 **{detected_country}**（可手动更改）")
    with col_t:
        project_type = st.selectbox("项目类型", ["Offshore", "Onshore"], index=0, key="su_ptype")
    is_offshore = project_type == "Offshore"
    profile = get_country_profile(country_name)

    eq_ratio = profile.typical_equity_ratio if profile else 0.30
    loan_rate = profile.typical_loan_rate if profile else 0.07
    loan_term = profile.typical_loan_term if profile else 15
    vat = profile.vat_rate if profile else 0.12
    cit = profile.corporate_income_tax_rate if profile else 0.25
    urban_tax = profile.urban_maintenance_tax_rate if profile else 0.0
    edu_sur = profile.education_surcharge_rate if profile else 0.0

    with st.expander("⚙️ 公共参数调整", expanded=False):
        col1, col2, col3 = st.columns(3)
        _su_tref = get_tariff_display(sel_country, "offshore" if is_offshore else "onshore")
        _su_tariff_def = round((_su_tref.low + _su_tref.high) / 2, 4) if _su_tref and _su_tref.high > 0 else 0.085
        tariff_override = col1.number_input("电价(含税 USD/kWh)", 0.001, 0.500, _su_tariff_def, step=0.001, format="%.4f", key="su_tariff")
        build_months = col2.number_input("建设期(月)", 6, 48, 24, key="su_build")
        oper_years = col3.number_input("运营期(年)", 15, 30, 25, key="su_oper")
        col4, col5, col6 = st.columns(3)
        eq_ratio = col4.number_input("资本金比例", 0.10, 0.80, eq_ratio, step=0.05, key="su_eq")
        loan_rate = col5.number_input("贷款利率", 0.01, 0.20, loan_rate, step=0.005, format="%.3f", key="su_lr")
        loan_term = col6.number_input("贷款期限(年)", 5, 25, loan_term, key="su_lt")
        if _su_tref:
            st.caption(f"📋 电价参考: {_su_tref.low:.4f}~{_su_tref.high:.4f} USD/kWh | {_su_tref.source} | {_su_tref.mechanism}")

    # ── 运维算法选择（Smart Upload模式 — 使用公共函数） ──
    with st.expander("🔧 运维算法选择", expanded=False):
        _su_om = _render_om_params(
            key_prefix="su", is_offshore=is_offshore, profile=profile,
            operation_years=oper_years, country_display=sel_country, compact=True,
        )

    with st.expander("📝 项目概况 (选填，导出PPT时展示)", expanded=False):
        su_pi_loc = st.text_input("项目地点", value="", placeholder="如：越南河静省 (Ha Tinh Province)", key="su_pi_loc")
        su_pi_tariff = st.text_input("电价来源", value="", placeholder="如：Decision 1508/QĐ-BCT (2025.05.30)", key="su_pi_tariff")
        su_pi_desc = st.text_area("项目说明/备注", value="", placeholder="如：近海风电(Nearshore)，2机型对比方案", height=68, key="su_pi_desc")

    if profile and profile.has_wind_tax_incentive:
        tax_holiday = profile.income_tax_holiday
    else:
        tax_holiday = (1, 3, 0.0, 4, 6, cit / 2.0)

    # ── 项目分组前缀 ──
    st.markdown("---")
    su_project_prefix = st.text_input(
        "📁 项目组名（Group）",
        value="",
        placeholder="如：Ha Tinh、Laguna、Frontera — 留空则归入 Smart 组",
        key="su_project_prefix",
        help="保存后所有方案将以「[组名] - [机型]」命名，自动归入同一个项目组。",
    )

    # ── Step 1: 构建可编辑表格 ──
    st.markdown("### 📊 参数确认 — 请核对并修改")
    st.caption("识别结果已填入下表，您可以**直接编辑**每个单元格。标红项为必须修正的错误，橙色为风险提示。")

    edit_rows = []
    for i, v in enumerate(variants):
        edit_rows.append({
            "方案": v.get("wtg_type", f"方案{i+1}"),
            "台数": _safe_float(v.get("units", None)),
            "单机MW": _safe_float(v.get("turbine_mw", None)),
            "P90(h)": _safe_float(v.get("p90_hours", None)),
            "电价(USD/kWh)": _safe_float(v.get("tariff_usd", None)),
            "TSI($/kW)": _safe_float(v.get("tsi_per_kw", None)),
            "BOP($/kW)": _safe_float(v.get("bop_per_kw", None)),
            "CAPEX($/kW)": _safe_float(v.get("capex_per_kw", None)),
        })

    edit_df = pd.DataFrame(edit_rows)
    col_config = {
        "方案": st.column_config.TextColumn("方案/机型", width="medium"),
        "台数": st.column_config.NumberColumn("台数", min_value=1, max_value=500, step=1, format="%d"),
        "单机MW": st.column_config.NumberColumn("单机MW", min_value=0.5, max_value=30.0, step=0.5, format="%.1f"),
        "P90(h)": st.column_config.NumberColumn("P90(h)", min_value=500, max_value=6000, step=10, format="%d"),
        "电价(USD/kWh)": st.column_config.NumberColumn("电价(USD/kWh)", min_value=0.0, max_value=0.500, step=0.001, format="%.4f"),
        "TSI($/kW)": st.column_config.NumberColumn("TSI($/kW)", min_value=0, max_value=5000, step=10, format="%.0f"),
        "BOP($/kW)": st.column_config.NumberColumn("BOP($/kW)", min_value=0, max_value=5000, step=10, format="%.0f"),
        "CAPEX($/kW)": st.column_config.NumberColumn("CAPEX($/kW)", min_value=0, max_value=8000, step=10, format="%.0f"),
    }

    edited_df = st.data_editor(
        edit_df,
        column_config=col_config,
        use_container_width=True,
        hide_index=True,
        num_rows="dynamic",
        key="su_editor",
    )

    # ── Step 2: 数据校验 ──
    errors, warnings = _validate_smart_params(edited_df, tariff_override, is_offshore)

    has_errors = len(errors) > 0
    has_warnings = len(warnings) > 0

    if has_errors:
        st.markdown("#### :red[错误 — 必须修正后才能计算]")
        for e in errors:
            st.error(e, icon="🚫")
    if has_warnings:
        st.markdown("#### :orange[风险提示 — 请确认后继续]")
        for w in warnings:
            st.warning(w, icon="⚠️")

    if not has_errors and not has_warnings:
        st.success("所有参数校验通过，可以计算！", icon="✅")

    # ── Step 3: 计算按钮 — 分级控制 ──
    can_calc = False
    if has_errors:
        st.button("🚀 计算全部方案", type="primary", key="su_calc", disabled=True,
                  help="存在必须修正的错误，请先编辑上方表格")
    elif has_warnings:
        confirmed = st.checkbox("我已确认上述风险项，继续计算", key="su_risk_confirm")
        if not confirmed:
            st.button("🚀 计算全部方案", type="primary", key="su_calc_disabled", disabled=True,
                      help="请先勾选确认风险项")
        else:
            if st.button("🚀 确认并计算全部方案", type="primary", key="su_calc"):
                can_calc = True
    else:
        if st.button("🚀 计算全部方案", type="primary", key="su_calc"):
            can_calc = True

    if can_calc:
        all_results = []
        valid_df = edited_df.dropna(subset=["CAPEX($/kW)"], how="all")
        valid_df = valid_df[valid_df["方案"].notna() & (valid_df["方案"] != "")]
        if valid_df.empty:
            st.error("没有有效的方案数据可以计算")
            return None
        for idx, row in valid_df.iterrows():
            wtg = str(row["方案"]) if pd.notna(row["方案"]) else f"方案{idx+1}"
            units = int(row["台数"]) if pd.notna(row["台数"]) and row["台数"] > 0 else 40
            mw = float(row["单机MW"]) if pd.notna(row["单机MW"]) and row["单机MW"] > 0 else 10.0
            p90 = int(row["P90(h)"]) if pd.notna(row["P90(h)"]) and row["P90(h)"] > 0 else 2500
            tariff_val = float(row["电价(USD/kWh)"]) if pd.notna(row["电价(USD/kWh)"]) and row["电价(USD/kWh)"] > 0 else tariff_override
            tsi = float(row["TSI($/kW)"]) if pd.notna(row["TSI($/kW)"]) else 0
            bop = float(row["BOP($/kW)"]) if pd.notna(row["BOP($/kW)"]) else 0
            capex = float(row["CAPEX($/kW)"]) if pd.notna(row["CAPEX($/kW)"]) and row["CAPEX($/kW)"] > 0 else (tsi + bop)

            _prefix = st.session_state.get("su_project_prefix", "").strip() or "Smart"
            _proj_name = f"{_prefix} - {wtg}"
            basic = BasicInfo(
                project_name=_proj_name, project_type="offshore" if is_offshore else "onshore",
                country=country_name, num_turbines=units, turbine_capacity_mw=mw,
                full_load_hours=p90, loss_rate=0.0, construction_months=build_months,
            )
            investment = InvestmentData(unit_static_investment=capex, working_capital_per_kw=4.0, deductible_vat_ratio=0.0)
            financing = FinancingTerms(
                equity_ratio=eq_ratio, long_term_loan_rate=loan_rate, loan_term_years=loan_term,
                working_capital_loan_rate=loan_rate, working_capital_equity_ratio=eq_ratio,
            )
            offshore_extra = OffshoreExtraCost(requires_sov=False, sea_area_usage_fee=43.0) if is_offshore else None
            operational = OperationalCost(
                om_method=_su_om["om_method"],
                staff_count=35 if is_offshore else 15, salary_per_person=3.5 if is_offshore else 1.0,
                welfare_rate=0.40, insurance_rate=0.0035, depreciation_years=20, residual_rate=0.0,
                operation_years=oper_years,
                warranty=_su_om["warranty"], post_warranty=_su_om["post_warranty"],
                base_om_per_kw=_su_om["base_om_per_kw"],
                om_escalation_rate=_su_om["om_escalation_rate"],
                capex_om_percentage=_su_om["capex_om_percentage"],
                capex_om_escalation=_su_om["capex_om_escalation"],
                contract_om_periods=_su_om["contract_om_periods"],
                offshore_extra=offshore_extra,
            )
            tax_financial = TaxAndFinancial(
                tariff_with_tax=tariff_val, vat_rate=vat, vat_refund_rate=0.0,
                income_tax_rate=cit, income_tax_holiday=tax_holiday,
                urban_maintenance_tax_rate=urban_tax, education_surcharge_rate=edu_sur,
                discount_rate=0.08,
            )
            inp = WindFarmFinancialInputs(basic=basic, investment=investment, financing=financing,
                                          operational=operational, tax_financial=tax_financial)
            res = calculate(inp)
            all_results.append((wtg, row, inp, res))

            save_project(_proj_name, inp, res, group=_prefix)

        st.success(f"已计算并保存 {len(all_results)} 个方案！")

        result_rows = []
        for wtg, row, inp, res in all_results:
            cap = int(row["台数"] or 40) * float(row["单机MW"] or 10)
            result_rows.append({
                "机型": wtg,
                "容量(MW)": f"{cap:.0f}",
                "CAPEX($/kW)": f"{row['CAPEX($/kW)'] or 0:.0f}",
                "全投IRR税前": f"{res.project_irr_before_tax*100:.2f}%",
                "全投IRR税后": f"{res.project_irr_after_tax*100:.2f}%",
                "资本金IRR": f"{res.equity_irr*100:.2f}%",
                "LCOE($/kWh)": f"{res.lcoe:.5f}",
                "NPV税后(M$)": f"{res.project_npv_after_tax/1e6:.1f}",
                "回收期(年)": f"{res.payback_after_tax:.1f}",
            })
        st.markdown("### 📈 计算结果对比")
        st.dataframe(pd.DataFrame(result_rows), use_container_width=True, hide_index=True)

        st.info("方案已自动保存。可切换到「项目对比」标签页查看详细对比图表，或导出对比 PPT。")

    return None


# ════════════════════════════════════════════════════════════════════════════
# 侧边栏：快速模式 (10 个参数即可)
# ════════════════════════════════════════════════════════════════════════════

def sidebar_inputs_quick() -> WindFarmFinancialInputs:
    """快速模式：只需截图级别的数据即可完成评估"""

    st.sidebar.markdown("## Quick Mode")
    st.sidebar.caption("Only 10 params needed. Defaults by country & type.")

    countries = list_countries()
    country_options = {f"{cn} ({en})": en for en, cn in countries}
    selected_display = st.sidebar.selectbox("Country", list(country_options.keys()), index=0, key="q_country")
    country_name = country_options[selected_display]
    profile = get_country_profile(country_name)

    project_type = st.sidebar.radio("Type", ["Onshore", "Offshore"], horizontal=True, key="q_type")
    is_offshore = project_type == "Offshore"

    # 海上/陆上不同的默认值
    _def = {
        "units":     (38,    28),
        "mw":        (9.0,   15.0),
        "p90":       (3429,  3200),
        "tariff":    (0.098, 0.098),
        "tsi":       (617.0, 900.0),
        "bop":       (433.0, 700.0),
        "build_m":   (18,    30),
        "oper_y":    (25,    25),
        "loss":      (0.03,  0.04),
        "staff":     (15,    35),
        "salary":    (1.0,   3.5),
        "insurance": (0.0025, 0.0035),
        "w_mat":     (3.0,   5.0),
        "w_other":   (5.0,   8.0),
        "pw_mat":    (4.0,   6.0),
        "pw_other":  (6.0,   10.0),
    }
    def D(key):
        return _def[key][1 if is_offshore else 0]

    st.sidebar.markdown("---")
    st.sidebar.markdown("### Basic")
    project_name = st.sidebar.text_input("Project Name", value="Quick Project", key="q_name")
    c1, c2 = st.sidebar.columns(2)
    num_turbines = c1.number_input("Units", 1, 500, D("units"), key="q_units")
    turbine_mw = c2.number_input("MW/unit", 1.0, 30.0, D("mw"), step=0.5, key="q_mw")
    full_load_hours = st.sidebar.number_input("P90 Hours (h/yr)", 1000, 5000, D("p90"), step=10, key="q_p90")
    tariff = st.sidebar.number_input("Tariff incl. tax (USD/kWh)", 0.001, 0.500, D("tariff"),
                                     step=0.001, format="%.4f", key="q_tariff")
    _q_tref = get_tariff_display(country_name, "offshore" if is_offshore else "onshore")
    if _q_tref:
        st.sidebar.caption(f"📋 参考: {_q_tref.low:.4f}~{_q_tref.high:.4f} USD/kWh ({_q_tref.source})")

    st.sidebar.markdown("---")
    st.sidebar.markdown("### Investment")
    tsi_per_kw = st.sidebar.number_input("TSI (USD/kW)", 100.0, 5000.0, D("tsi"), step=10.0, key="q_tsi")
    bop_per_kw = st.sidebar.number_input("BOP (USD/kW)", 50.0, 3000.0, D("bop"), step=10.0, key="q_bop")
    total_per_kw = tsi_per_kw + bop_per_kw
    capacity_mw = num_turbines * turbine_mw
    total_invest_m = total_per_kw * capacity_mw * 1000 / 1e6
    st.sidebar.info(f"**Total: {total_per_kw:,.0f} USD/kW | {total_invest_m:,.0f} M$**")

    st.sidebar.markdown("---")
    st.sidebar.markdown("### Timeline")
    c3, c4 = st.sidebar.columns(2)
    construction_months = c3.number_input("Build (months)", 6, 48, D("build_m"), key="q_build")
    operation_years = c4.number_input("Operate (years)", 15, 30, D("oper_y"), key="q_oper")

    # 从国家配置自动填充
    eq_ratio = profile.typical_equity_ratio if profile else 0.30
    loan_rate = profile.typical_loan_rate if profile else 0.07
    loan_term = profile.typical_loan_term if profile else 15
    vat = profile.vat_rate if profile else 0.12
    cit = profile.corporate_income_tax_rate if profile else 0.25
    urban_tax = profile.urban_maintenance_tax_rate if profile else 0.0
    edu_sur = profile.education_surcharge_rate if profile else 0.0

    if profile and profile.has_wind_tax_incentive:
        tax_holiday = (1, 7, 0.0, 8, 14, cit * 0.4)
    else:
        tax_holiday = (1, 3, 0.0, 4, 6, cit / 2.0)

    # ── 运维算法选择（Quick模式 — 使用公共函数） ──
    _q_om = _render_om_params(
        key_prefix="q", is_offshore=is_offshore, profile=profile,
        operation_years=operation_years,
        total_investment_usd=total_per_kw * num_turbines * turbine_mw * 1000,
        country_display=selected_display, compact=True, use_sidebar=True,
    )

    # 显示自动填充的默认值
    with st.sidebar.expander("Auto-filled defaults (view only)", expanded=False):
        st.markdown(f"""
- **Equity ratio**: {eq_ratio:.0%}
- **Loan rate**: {loan_rate:.1%} x {loan_term}yr
- **VAT**: {vat:.0%} | **CIT**: {cit:.0%}
- **Loss rate**: {D('loss'):.0%}
- **Staff**: {D('staff')} | **Insurance**: {D('insurance'):.2%}
- **Type**: {'Offshore' if is_offshore else 'Onshore'} defaults
        """)
        st.caption("Switch to Detailed mode to edit these.")

    onshore_detail = None
    offshore_extra = None

    if is_offshore:
        offshore_extra = OffshoreExtraCost(
            requires_sov=False, sov_annual_cost=0.0,
            sea_area_usage_fee=43.0, storage_rental=0.0,
            decommissioning_rate=0.02,
        )
    else:
        onshore_detail = OnshoreInvestment(
            equipment_and_installation=tsi_per_kw,
            civil_works=bop_per_kw * 0.85,
            construction_auxiliary=0.0,
            other_costs=bop_per_kw * 0.15,
            contingency_rate=0.0,
            storage_cost=0.0,
            grid_connection_cost=0.0,
        )

    basic = BasicInfo(
        project_name=project_name,
        project_type="offshore" if is_offshore else "onshore",
        country=country_name,
        num_turbines=num_turbines,
        turbine_capacity_mw=turbine_mw,
        full_load_hours=full_load_hours,
        loss_rate=D("loss"),
        construction_months=construction_months,
    )
    investment = InvestmentData(
        unit_static_investment=total_per_kw,
        working_capital_per_kw=4.0,
        deductible_vat_ratio=0.0,
        onshore_detail=onshore_detail,
        offshore_detail=None,
    )
    financing = FinancingTerms(
        equity_ratio=eq_ratio,
        long_term_loan_rate=loan_rate,
        loan_term_years=loan_term,
        working_capital_loan_rate=loan_rate,
        working_capital_equity_ratio=eq_ratio,
    )
    operational = OperationalCost(
        om_method=_q_om["om_method"],
        staff_count=D("staff"), salary_per_person=D("salary"), welfare_rate=0.40,
        insurance_rate=D("insurance"), depreciation_years=20, residual_rate=0.0,
        operation_years=operation_years,
        warranty=_q_om["warranty"], post_warranty=_q_om["post_warranty"],
        base_om_per_kw=_q_om["base_om_per_kw"],
        om_escalation_rate=_q_om["om_escalation_rate"],
        capex_om_percentage=_q_om["capex_om_percentage"],
        capex_om_escalation=_q_om["capex_om_escalation"],
        contract_om_periods=_q_om["contract_om_periods"],
        offshore_extra=offshore_extra,
    )
    tax_financial = TaxAndFinancial(
        tariff_with_tax=tariff, vat_rate=vat, vat_refund_rate=0.0,
        income_tax_rate=cit, income_tax_holiday=tax_holiday,
        urban_maintenance_tax_rate=urban_tax, education_surcharge_rate=edu_sur,
        discount_rate=0.08,
    )
    return WindFarmFinancialInputs(
        basic=basic, investment=investment, financing=financing,
        operational=operational, tax_financial=tax_financial,
    )


# ════════════════════════════════════════════════════════════════════════════
# 侧边栏：完整分项参数编辑
# ════════════════════════════════════════════════════════════════════════════

def sidebar_inputs() -> WindFarmFinancialInputs:
    """完整的侧边栏参数面板，含所有分项编辑"""

    st.sidebar.markdown("## ⚙️ 项目参数设定")

    # ──── 国家 ────
    countries = list_countries()
    country_options = {f"{cn} ({en})": en for en, cn in countries}
    selected_display = st.sidebar.selectbox(
        "🌏 国家/地区", list(country_options.keys()), index=0
    )
    country_name = country_options[selected_display]
    profile = get_country_profile(country_name)

    # ──── 项目类型 ────
    project_type = st.sidebar.radio(
        "项目类型", ["陆上风电 (Onshore)", "海上风电 (Offshore)"], horizontal=True
    )
    is_offshore = "Offshore" in project_type

    _sens_help = st.sidebar.checkbox("📊 显示参数说明与敏感性", value=False, key="show_sens_top")

    st.sidebar.markdown("---")

    # ═══════════════════ 1. 基本信息 ═══════════════════
    st.sidebar.markdown("### 📋 基本信息")
    project_name = st.sidebar.text_input("项目名称", value="Demo Wind Farm")
    c1, c2 = st.sidebar.columns(2)
    num_turbines = c1.number_input("机组台数", 1, 500, 28 if is_offshore else 16)
    turbine_mw = c2.number_input("单机容量(MW)", 1.0, 30.0, 18.0 if is_offshore else 6.25, step=0.5)
    full_load_hours = st.sidebar.number_input("满负荷小时数 (h)", 1000, 5000, 3138 if is_offshore else 2523, step=10)
    if _sens_help:
        st.sidebar.caption("⚡⚡ **最高敏感** | ±100h ≈ IRR±0.3~0.8pp。是决定项目盈利的核心参数，建议用P50/P75/P90区分风险场景。")
    loss_rate = st.sidebar.number_input("综合线损率 (%)", 0.0, 10.0, 3.0, step=0.5) / 100.0
    if _sens_help:
        st.sidebar.caption("🔵 **中敏感** | ±1% ≈ IRR±0.1~0.2pp。含厂用电+线路损耗，陆上2~4%，海上2~5%。")
    construction_months = st.sidebar.number_input("建设期 (月)", 6, 36, 24 if is_offshore else 12, step=1)
    if _sens_help:
        st.sidebar.caption("🔵 **中敏感** | 建设期延长→建设期利息增加→IRR下降。每多6个月约IRR-0.1~0.3pp。")

    st.sidebar.markdown("---")

    # ═══════════════════ 2. 投资造价分项 ═══════════════════
    st.sidebar.markdown("### 💰 投资造价")

    onshore_detail = None
    offshore_detail = None

    if is_offshore:
        # ---- 海上 EPC 明细 ----
        with st.sidebar.expander("🔩 OEM 成本 (风机+塔筒)", expanded=False):
            oem_turbine_price = st.number_input("风机售价 (USD/kW)", 100.0, 3000.0, 680.0, step=10.0, key="oem_tp")
            oem_tower_weight = st.number_input("塔筒重量+内附件 (t/台)", 50.0, 2000.0, 600.0, step=10.0, key="oem_tw")
            oem_tower_price = st.number_input("塔筒单价 (USD/吨)", 500.0, 5000.0, 1200.0, step=50.0, key="oem_tup")

        with st.sidebar.expander("🚢 安装与运输", expanded=False):
            inst_install = st.number_input("安装费 (USD/台)", 0.0, 5e6, 800_000.0, step=50_000.0, key="inst_i")
            inst_ocean = st.number_input("海运费 (USD/台)", 0.0, 5e6, 400_000.0, step=50_000.0, key="inst_o")
            inst_inland = st.number_input("陆运费 (USD/台)", 0.0, 2e6, 50_000.0, step=10_000.0, key="inst_l")

        with st.sidebar.expander("🏗️ 基础工程", expanded=False):
            fnd_price_per_ton = st.number_input("基础造价 (USD/吨)", 500.0, 10_000.0, 3500.0, step=100.0, key="fnd_p")
            fnd_tons = st.number_input("基础重量 (t/台)", 100.0, 5000.0, 1500.0, step=50.0, key="fnd_t")
            fnd_install = st.number_input("基础安装费 (USD/台)", 0.0, 5e6, 500_000.0, step=50_000.0, key="fnd_i")

        with st.sidebar.expander("⚡ BOP 分项工程 (USD/kW)", expanded=False):
            bop_aux = st.number_input("施工辅助工程", 0.0, 500.0, 20.0, step=1.0, key="bop1")
            bop_coll_eq = st.number_input("集电线路设备", 0.0, 500.0, 40.0, step=1.0, key="bop2")
            bop_subcable = st.number_input("海缆设备", 0.0, 500.0, 80.0, step=1.0, key="bop3")
            bop_offshore_sub = st.number_input("海上升压站设备", 0.0, 500.0, 60.0, step=1.0, key="bop4")
            bop_ctrl = st.number_input("集控中心", 0.0, 200.0, 15.0, step=1.0, key="bop5")
            bop_other_eq = st.number_input("其他设备", 0.0, 200.0, 10.0, step=1.0, key="bop6")
            bop_coll_civ = st.number_input("集电线路工程", 0.0, 200.0, 20.0, step=1.0, key="bop7")
            bop_landing = st.number_input("登陆电缆工程", 0.0, 200.0, 30.0, step=1.0, key="bop8")
            bop_sub_civ = st.number_input("海上升压站工程", 0.0, 200.0, 25.0, step=1.0, key="bop9")
            bop_ctrl_civ = st.number_input("集控中心工程", 0.0, 200.0, 10.0, step=1.0, key="bop10")
            bop_transp = st.number_input("交通工程", 0.0, 200.0, 5.0, step=1.0, key="bop11")
            bop_other_civ = st.number_input("其他工程", 0.0, 200.0, 5.0, step=1.0, key="bop12")

        oem = OEMCost(oem_turbine_price, oem_tower_weight, oem_tower_price)
        installation = InstallationCost(inst_install, inst_ocean, inst_inland)
        foundation = FoundationCost(fnd_price_per_ton, fnd_tons, fnd_install)
        bop = BOPCost(
            bop_aux, bop_coll_eq, bop_subcable, bop_offshore_sub, bop_ctrl, bop_other_eq,
            bop_coll_civ, bop_landing, bop_sub_civ, bop_ctrl_civ, bop_transp, bop_other_civ,
        )
        offshore_detail = OffshoreEPCBreakdown(
            oem=oem, installation=installation, foundation=foundation, bop=bop,
            num_turbines=num_turbines, turbine_capacity_mw=turbine_mw,
        )
        auto_epc = offshore_detail.total_epc_per_kw
        st.sidebar.info(f"EPC 明细合计: **{auto_epc:,.1f} USD/kW**")
        use_detail = st.sidebar.checkbox("使用 EPC 明细计算投资", value=True, key="use_epc")
        if use_detail:
            unit_investment = auto_epc
        else:
            unit_investment = st.sidebar.number_input("手动输入 (USD/kW)", 200.0, 5000.0, 1833.8, step=10.0, key="inv_manual_off")

    else:
        # ---- 陆上投资明细 ----
        with st.sidebar.expander("📦 陆上投资明细 (USD/kW)", expanded=False):
            on_equip = st.number_input("设备及安装工程", 0.0, 3000.0, 550.0, step=10.0, key="on1")
            on_civil = st.number_input("建筑工程", 0.0, 1000.0, 100.0, step=5.0, key="on2")
            on_aux = st.number_input("施工辅助工程", 0.0, 200.0, 15.0, step=1.0, key="on3")
            on_other = st.number_input("其他费用", 0.0, 500.0, 60.0, step=5.0, key="on4")
            on_contingency = st.number_input("基本预备费率 (%)", 0.0, 10.0, 2.0, step=0.5, key="on5") / 100.0
            on_storage = st.number_input("储能工程", 0.0, 500.0, 30.0, step=5.0, key="on6")
            on_grid = st.number_input("送出线路/电网接入", 0.0, 500.0, 50.0, step=5.0, key="on7")

        onshore_detail = OnshoreInvestment(
            on_equip, on_civil, on_aux, on_other, on_contingency, on_storage, on_grid,
        )
        auto_onshore = onshore_detail.total_per_kw
        st.sidebar.info(f"陆上明细合计: **{auto_onshore:,.1f} USD/kW**")
        use_detail = st.sidebar.checkbox("使用陆上明细计算投资", value=True, key="use_on")
        if use_detail:
            unit_investment = auto_onshore
        else:
            unit_investment = st.sidebar.number_input("手动输入 (USD/kW)", 200.0, 5000.0, 816.9, step=10.0, key="inv_manual_on")

    working_capital_per_kw = st.sidebar.number_input("流动资金 (USD/kW)", 0.0, 50.0, 4.2, step=0.5)
    if _sens_help:
        st.sidebar.caption("⚡⚡ **最高敏感** | 单位投资±100$/kW ≈ IRR±0.5~1.5pp。是仅次于电价和发电量的第三大敏感因子。")

    st.sidebar.markdown("---")

    # ═══════════════════ 3. 融资条件 ═══════════════════
    st.sidebar.markdown("### 🏦 融资条件")
    default_eq = profile.typical_equity_ratio * 100 if profile else 25.0
    default_rate = profile.typical_loan_rate * 100 if profile else 3.25
    default_term = profile.typical_loan_term if profile else 15

    equity_ratio = st.sidebar.number_input("资本金比例 (%)", 10.0, 50.0, default_eq, step=1.0) / 100.0
    if _sens_help:
        st.sidebar.caption("⚡ **高敏感** | 资本金↑10% ≈ 资本金IRR↓2~4pp（杠杆效应减弱），但全投资IRR不变。影响资本金IRR和融资结构。")
    loan_rate = st.sidebar.number_input("贷款年利率 (%)", 0.5, 15.0, default_rate, step=0.25, format="%.2f") / 100.0
    if _sens_help:
        st.sidebar.caption(f"⚡ **高敏感** | 利率±1% ≈ IRR±0.3~0.7pp。{selected_display}当前: {default_rate:.2f}%")
    loan_term = st.sidebar.number_input("贷款年限", 5, 25, default_term, step=1)
    if _sens_help:
        st.sidebar.caption("🔵 **中敏感** | 延长贷款=前期还款压力小=资本金IRR提升，但总利息增加。")

    with st.sidebar.expander("🔄 流动资金贷款", expanded=False):
        wc_loan_rate = st.number_input("流动资金贷款利率 (%)", 0.0, 15.0, 3.25, step=0.25, key="wclr") / 100.0
        wc_equity_ratio = st.number_input("流动资金自有资金比例 (%)", 0.0, 100.0, 30.0, step=5.0, key="wcer") / 100.0
        if _sens_help:
            st.caption("⚪ **低敏感** | 流动资金占总投资<1%，对IRR影响极小（<0.05pp）。")

    st.sidebar.markdown("---")

    # ═══════════════════ 4. 税费与电价 ═══════════════════
    st.sidebar.markdown("### 📊 税费与电价")
    default_vat = profile.vat_rate * 100 if profile else 13.0
    default_cit = profile.corporate_income_tax_rate * 100 if profile else 25.0
    _tariff_range = profile.offshore_tariff_range if is_offshore else profile.onshore_tariff_range
    if _tariff_range and _tariff_range[0] > 0:
        default_tariff = round((_tariff_range[0] + _tariff_range[1]) / 2, 4)
    else:
        default_tariff = 0.0638 if is_offshore else 0.0434

    tariff = st.sidebar.number_input("含税电价 (USD/kWh)", 0.001, 0.500, default_tariff, step=0.001, format="%.4f")
    _tariff_ref = get_tariff_display(country_name, "offshore" if is_offshore else "onshore")
    _tariff_policy = _tariff_ref.policy_text if _tariff_ref else ""
    _tariff_mech = _tariff_ref.mechanism if _tariff_ref else ""
    _tariff_src = _tariff_ref.source if _tariff_ref else ""
    _ref_lo = _tariff_ref.low if _tariff_ref else _tariff_range[0]
    _ref_hi = _tariff_ref.high if _tariff_ref else _tariff_range[1]
    if _tariff_ref or _tariff_policy:
        with st.sidebar.expander(f"📋 {selected_display} {'海上' if is_offshore else '陆上'}电价参考", expanded=False):
            if _tariff_mech:
                st.markdown(f"**Mechanism**: {_tariff_mech}")
            st.markdown(f"**Range**: {_ref_lo:.4f} ~ {_ref_hi:.4f} USD/kWh")
            if _tariff_src:
                st.caption(f"数据来源: {_tariff_src}")
            if _tariff_policy:
                st.markdown(f"**Policy**: {_tariff_policy}")
    if _sens_help:
        st.sidebar.caption(f"⚡⚡ **最高敏感** | ±0.01$/kWh ≈ IRR±1~2pp。{selected_display}参考: {'海上' if is_offshore else '陆上'} {_ref_lo:.4f}~{_ref_hi:.4f} USD/kWh")
    vat_rate = st.sidebar.number_input("增值税率 (%)", 0.0, 20.0, default_vat, step=1.0) / 100.0
    vat_refund = st.sidebar.number_input("即征即退比例 (%)", 0.0, 100.0, 50.0 if country_name == "China" else 0.0, step=5.0) / 100.0
    income_tax_rate = st.sidebar.number_input("所得税率 (%)", 0.0, 35.0, default_cit, step=1.0) / 100.0
    if _sens_help:
        st.sidebar.caption(f"🔵 **中敏感** | 所得税±5% ≈ IRR税后±0.2~0.4pp。{selected_display}: {default_cit:.0f}%{'，有风电优惠' if profile and profile.has_wind_tax_incentive else ''}")
    discount_rate = st.sidebar.number_input("基准折现率 (%)", 3.0, 15.0, 8.0, step=0.5) / 100.0
    if _sens_help:
        st.sidebar.caption("🔵 **中敏感** | 折现率不影响IRR，只影响NPV。折现率↑=NPV↓。8%为中国标准，海外项目常用WACC(6~10%)。")

    with st.sidebar.expander("📑 税费附加 & 所得税优惠", expanded=False):
        urban_tax = st.number_input("城市维护建设税率 (%)", 0.0, 10.0,
                                     (profile.urban_maintenance_tax_rate if profile else 0.05) * 100,
                                     step=1.0, key="utax") / 100.0
        edu_surcharge = st.number_input("教育费附加率 (%)", 0.0, 10.0,
                                         (profile.education_surcharge_rate if profile else 0.05) * 100,
                                         step=1.0, key="edu") / 100.0

        st.markdown("**所得税优惠政策**")
        if profile and profile.has_wind_tax_incentive:
            st.caption(f"当前国家优惠: {profile.tax_incentive_description}")
        c1, c2, c3 = st.columns(3)
        exempt_start = c1.number_input("免征起始年", 1, 25, 1, key="exs")
        exempt_end = c2.number_input("免征结束年", 1, 25, 3, key="exe")
        exempt_rate = c3.number_input("免征期税率", 0.0, 0.5, 0.0, step=0.01, key="exr")
        c4, c5, c6 = st.columns(3)
        half_start = c4.number_input("减半起始年", 1, 25, 4, key="hfs")
        half_end = c5.number_input("减半结束年", 1, 25, 6, key="hfe")
        half_rate = c6.number_input("减半期税率", 0.0, 0.5, income_tax_rate / 2.0, step=0.01, key="hfr")
        tax_holiday = (exempt_start, exempt_end, exempt_rate, half_start, half_end, half_rate)

    st.sidebar.markdown("---")

    # ═══════════════════ 5. 运营成本分项 ═══════════════════
    st.sidebar.markdown("### 🔧 运营成本")

    operation_years = st.sidebar.number_input("运营期 (年)", 15, 35, 25 if is_offshore else 20, step=1)
    if _sens_help:
        st.sidebar.caption("⚡ **高敏感** | 延长运营期可提升IRR 0.3~1.0pp，但后期运维成本递增会削弱效果。海上25年、陆上20年为行业惯例。")

    depreciation_years = st.sidebar.number_input("折旧年限 (年)", 10, 25, 20, step=1)
    if _sens_help:
        st.sidebar.caption("🔵 **中敏感** | 影响年折旧额→利润→所得税。缩短折旧年限=前期少交税=IRR微升，但对银行还款计划无影响。")

    residual_rate = st.sidebar.number_input("残值率 (%)", 0.0, 10.0, 0.0 if is_offshore else 5.0, step=0.5) / 100.0
    if _sens_help:
        st.sidebar.caption("⚪ **低敏感** | 仅影响折旧基数和末年回收。海上通常0%（拆除成本抵消），陆上5%。")

    with st.sidebar.expander("👷 人员与保险", expanded=False):
        staff_count = st.number_input("定员人数", 1, 200, 35 if is_offshore else 18, key="staff")
        salary_per_person = st.number_input("人均年薪 (万USD)", 0.5, 20.0, 3.52 if is_offshore else 1.11, step=0.1, key="sal")
        welfare_rate = st.number_input("福利系数", 0.0, 2.0, 0.60, step=0.05, key="welf")
        insurance_rate = st.number_input("保险费率 (%)", 0.0, 1.0, 0.35 if is_offshore else 0.25, step=0.05, key="ins") / 100.0
        if _sens_help:
            st.caption("⚪ 人员和保险为固定成本，占运维总额10~25%。对IRR影响有限（<0.1pp），但影响LCOE约1~3厘/kWh。")

    # ── 运维算法选择（Detailed模式 — 使用公共函数） ──
    _d_om = _render_om_params(
        key_prefix="detail", is_offshore=is_offshore, profile=profile,
        operation_years=operation_years,
        total_investment_usd=unit_investment * num_turbines * turbine_mw * 1000,
        country_display=selected_display, compact=False, use_sidebar=True,
    )

    # ──── 海上专项 ────
    offshore_extra = None
    if is_offshore:
        with st.sidebar.expander("🚢 海上专项费用", expanded=False):
            requires_sov = st.checkbox("需要 SOV 运维船", value=False, key="sov")
            sov_cost = st.number_input("SOV 年费 (万USD)", 0.0, 500.0, 0.0, key="sov_c") if requires_sov else 0.0
            sea_area_fee = st.number_input("海域使用金 (万USD/年)", 0.0, 500.0, 43.2, step=1.0, key="sea")
            storage_rental = st.number_input("储能租赁费 (万USD/年)", 0.0, 500.0, 0.0, step=1.0, key="stor_r")
            decomm_rate = st.number_input("拆除费率 (%)", 0.0, 10.0, 2.0, step=0.5, key="decomm") / 100.0
            offshore_extra = OffshoreExtraCost(
                requires_sov=requires_sov,
                sov_annual_cost=sov_cost,
                sea_area_usage_fee=sea_area_fee,
                storage_rental=storage_rental,
                decommissioning_rate=decomm_rate,
            )
            if _sens_help:
                st.caption("🔵 **中敏感** | SOV 150万$/年 ≈ IRR-0.3pp。海域金对IRR影响约0.05~0.15pp。是否需要SOV取决于离岸距离和水深。")

    # ═══════════════════ 组装 ═══════════════════
    basic = BasicInfo(
        project_name=project_name,
        project_type="offshore" if is_offshore else "onshore",
        country=country_name,
        num_turbines=num_turbines,
        turbine_capacity_mw=turbine_mw,
        full_load_hours=full_load_hours,
        loss_rate=loss_rate,
        construction_months=construction_months,
    )

    investment = InvestmentData(
        unit_static_investment=unit_investment,
        working_capital_per_kw=working_capital_per_kw,
        onshore_detail=onshore_detail,
        offshore_detail=offshore_detail,
    )

    financing = FinancingTerms(
        equity_ratio=equity_ratio,
        long_term_loan_rate=loan_rate,
        loan_term_years=loan_term,
        working_capital_loan_rate=wc_loan_rate,
        working_capital_equity_ratio=wc_equity_ratio,
    )

    operational = OperationalCost(
        om_method=_d_om["om_method"],
        staff_count=staff_count,
        salary_per_person=salary_per_person,
        welfare_rate=welfare_rate,
        insurance_rate=insurance_rate,
        depreciation_years=depreciation_years,
        residual_rate=residual_rate,
        operation_years=operation_years,
        warranty=_d_om["warranty"],
        post_warranty=_d_om["post_warranty"],
        base_om_per_kw=_d_om["base_om_per_kw"],
        om_escalation_rate=_d_om["om_escalation_rate"],
        capex_om_percentage=_d_om["capex_om_percentage"],
        capex_om_escalation=_d_om["capex_om_escalation"],
        contract_om_periods=_d_om["contract_om_periods"],
        offshore_extra=offshore_extra,
    )

    tax_financial = TaxAndFinancial(
        tariff_with_tax=tariff,
        vat_rate=vat_rate,
        vat_refund_rate=vat_refund,
        income_tax_rate=income_tax_rate,
        income_tax_holiday=tax_holiday,
        urban_maintenance_tax_rate=urban_tax,
        education_surcharge_rate=edu_surcharge,
        discount_rate=discount_rate,
    )

    return WindFarmFinancialInputs(
        basic=basic,
        investment=investment,
        financing=financing,
        operational=operational,
        tax_financial=tax_financial,
    )


# ════════════════════════════════════════════════════════════════════════════
# 图表函数（复用）
# ════════════════════════════════════════════════════════════════════════════

def plot_kpi_cards(result: CalculationResult):
    c1, c2, c3, c4, c5 = st.columns(5)
    c1.metric("全投资IRR(税后)", f"{result.project_irr_after_tax:.2%}")
    c2.metric("资本金IRR", f"{result.equity_irr:.2%}")
    c3.metric("LCOE", f"{result.lcoe:.5f} $/kWh")
    c4.metric("回收期(税后)", f"{result.payback_after_tax:.1f} 年")
    c5.metric("NPV(税后)", f"{result.project_npv_after_tax / 1e6:,.1f} M$")


def plot_investment_breakdown(inputs: WindFarmFinancialInputs):
    items = {"静态投资": inputs.total_static_investment,
             "建设期利息": inputs.construction_interest,
             "流动资金": inputs.working_capital}

    if inputs.investment.onshore_detail:
        od = inputs.investment.onshore_detail
        items = {
            "设备及安装": od.equipment_and_installation * inputs.capacity_kw,
            "建筑工程": od.civil_works * inputs.capacity_kw,
            "施工辅助": od.construction_auxiliary * inputs.capacity_kw,
            "其他费用": od.other_costs * inputs.capacity_kw,
            "储能": od.storage_cost * inputs.capacity_kw,
            "送出线路": od.grid_connection_cost * inputs.capacity_kw,
            "预备费": od.contingency * inputs.capacity_kw,
            "建设期利息": inputs.construction_interest,
        }
    elif inputs.investment.offshore_detail:
        od = inputs.investment.offshore_detail
        items = {
            "OEM(风机+塔筒)": od.oem_per_kw * inputs.capacity_kw,
            "安装与运输": od.installation_per_kw * inputs.capacity_kw,
            "基础工程": od.foundation_per_kw * inputs.capacity_kw,
            "BOP 合计": od.bop.total_bop_per_kw * inputs.capacity_kw,
            "建设期利息": inputs.construction_interest,
        }

    items = {k: v for k, v in items.items() if v > 0}
    fig = px.pie(
        pd.DataFrame({"项目": list(items.keys()), "金额": list(items.values())}),
        values="金额", names="项目",
        color_discrete_sequence=COLOR_PALETTE, hole=0.45,
    )
    fig.update_layout(title="投资构成", height=380, margin=dict(t=40, b=20))
    st.plotly_chart(fig, use_container_width=True)


def plot_cashflow_chart(result: CalculationResult):
    years, cf_after, cum = [], [], []
    running = 0.0
    for f in result.annual_flows:
        years.append(f"Y{f.year}")
        cf_after.append(f.project_net_cf_after_tax)
        running += f.project_net_cf_after_tax
        cum.append(running)

    fig = go.Figure()
    fig.add_trace(go.Bar(x=years, y=cf_after, name="税后净现金流", marker_color="#2E75B6", opacity=0.85))
    fig.add_trace(go.Scatter(x=years, y=cum, name="累计", line=dict(color="#C00000", width=2.5), yaxis="y2"))
    fig.update_layout(
        title="逐年全投资现金流", xaxis_title="年份",
        yaxis_title="净现金流 (USD)",
        yaxis2=dict(title="累计 (USD)", overlaying="y", side="right"),
        legend=dict(x=0.01, y=0.99), height=400, margin=dict(t=40, b=40),
    )
    st.plotly_chart(fig, use_container_width=True)


def plot_profit_chart(result: CalculationResult):
    op = [f for f in result.annual_flows if not f.is_construction]
    yrs = [f"Y{f.year}" for f in op]
    fig = go.Figure()
    fig.add_trace(go.Bar(x=yrs, y=[f.revenue for f in op], name="营业收入", marker_color="#548235"))
    fig.add_trace(go.Bar(x=yrs, y=[-f.total_opex for f in op], name="经营成本", marker_color="#BF8F00"))
    fig.add_trace(go.Bar(x=yrs, y=[-f.depreciation for f in op], name="折旧", marker_color="#7030A0"))
    fig.add_trace(go.Bar(x=yrs, y=[-(f.loan_interest + f.wc_loan_interest) for f in op], name="利息", marker_color="#C00000"))
    fig.add_trace(go.Scatter(x=yrs, y=[f.net_profit for f in op], name="净利润", line=dict(color="#1F4E79", width=2.5)))
    fig.update_layout(barmode="relative", title="利润结构", height=400, margin=dict(t=40, b=40))
    st.plotly_chart(fig, use_container_width=True)


def plot_sensitivity(inputs: WindFarmFinancialInputs):
    factors = {
        "电价": ("tariff_with_tax", inputs.tax_financial.tariff_with_tax),
        "投资": ("unit_static_investment", inputs.investment.unit_static_investment),
        "小时数": ("full_load_hours", inputs.basic.full_load_hours),
        "贷款利率": ("long_term_loan_rate", inputs.financing.long_term_loan_rate),
    }
    pct_range = [-20, -10, -5, 0, 5, 10, 20]
    heatmap, labels = [], []
    for name, (attr, base_val) in factors.items():
        labels.append(name)
        row = []
        for pct in pct_range:
            inp = copy.deepcopy(inputs)
            nv = base_val * (1.0 + pct / 100.0)
            if attr == "tariff_with_tax": inp.tax_financial.tariff_with_tax = nv
            elif attr == "unit_static_investment": inp.investment.unit_static_investment = nv
            elif attr == "full_load_hours": inp.basic.full_load_hours = max(500, int(round(nv)))
            elif attr == "long_term_loan_rate": inp.financing.long_term_loan_rate = max(0.001, nv)
            try:
                row.append(calculate(inp).project_irr_after_tax * 100)
            except Exception:
                row.append(0.0)
        heatmap.append(row)

    fig = go.Figure(data=go.Heatmap(
        z=heatmap, x=[f"{p:+d}%" for p in pct_range], y=labels,
        colorscale="RdYlGn",
        text=[[f"{v:.2f}%" for v in r] for r in heatmap],
        texttemplate="%{text}", textfont=dict(size=12),
        colorbar=dict(title="IRR(%)"),
    ))
    fig.update_layout(title="敏感性分析 — 全投资税后 IRR", height=350, margin=dict(t=40, b=40))
    st.plotly_chart(fig, use_container_width=True)


def generate_ppt_bytes(inputs: WindFarmFinancialInputs, result: CalculationResult,
                       project_info: Optional[dict] = None) -> bytes:
    """生成单项目经济性评估 PPT，返回二进制内容。
    project_info 可选，含 location/tariff_source/description 等自定义项目概况。"""
    from io import BytesIO
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK   = RGBColor(0x33, 0x33, 0x33)
    DARK    = RGBColor(0x1B, 0x2A, 0x4A)
    BLUE    = RGBColor(0x00, 0x6E, 0xB8)
    LBLUE   = RGBColor(0xD6, 0xEA, 0xF8)
    GREEN   = RGBColor(0x27, 0xAE, 0x60)
    LGREEN  = RGBColor(0xE8, 0xF8, 0xE8)
    GRAY    = RGBColor(0x7F, 0x8C, 0x8D)
    LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)
    ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
    RED_HL  = RGBColor(0xE7, 0x4C, 0x3C)

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    def _rect(l, t, w, h, fc=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.line.fill.background()
        if fc:
            s.fill.solid()
            s.fill.fore_color.rgb = fc
        return s

    def _txt(l, t, w, h, text, sz=12, c=BLACK, b=False, al=PP_ALIGN.LEFT):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = "Microsoft YaHei"
        p.alignment = al

    def _cell(tbl, r, c, text, sz=9, color=BLACK, bold=False,
              al=PP_ALIGN.CENTER, fc=None):
        cl = tbl.cell(r, c)
        cl.text = ""
        p = cl.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Microsoft YaHei"
        p.alignment = al
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fc:
            cl.fill.solid()
            cl.fill.fore_color.rgb = fc

    project_type_cn = "海上" if inputs.basic.project_type == "offshore" else "陆上"

    # ── 标题栏 ──
    _rect(Inches(0), Inches(0), SLIDE_W, Inches(0.9), fc=DARK)
    _txt(Inches(0.5), Inches(0.1), Inches(9), Inches(0.45),
         f"{inputs.basic.project_name} 经济性评估", sz=22, c=WHITE, b=True)
    _txt(Inches(0.5), Inches(0.5), Inches(10), Inches(0.35),
         f"{inputs.basic.country} · {project_type_cn}风电 · "
         f"{inputs.basic.num_turbines}×{inputs.basic.turbine_capacity_mw}MW = "
         f"{inputs.capacity_mw:.0f}MW · {inputs.operational.operation_years}年运营",
         sz=11, c=RGBColor(0xAA, 0xCC, 0xEE))

    logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mingyang_logo.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(11.3), Inches(0.1), height=Inches(0.7))

    # ── 左侧: 项目概况 ──
    Y0 = Inches(1.05)
    _rect(Inches(0.3), Y0, Inches(0.12), Inches(0.32), fc=BLUE)
    _txt(Inches(0.55), Y0, Inches(3), Inches(0.32), "项目概况", sz=14, c=DARK, b=True)

    tariff_usd = inputs.tax_financial.tariff_with_tax
    tariff_notax = inputs.tax_financial.tariff_without_tax

    pi = project_info or {}
    info_items = []
    if pi.get("location"):
        info_items.append(("项目地点", pi["location"]))
    info_items.append(("项目类型", f"{project_type_cn}风电 ({inputs.basic.project_type})"))
    info_items.append(("装机容量", f"{inputs.basic.num_turbines}台×{inputs.basic.turbine_capacity_mw}MW = {inputs.capacity_mw:.0f}MW"))
    info_items.append(("P90小时数", f"{inputs.basic.full_load_hours} h"))
    if pi.get("tariff_source"):
        info_items.append(("电价来源", pi["tariff_source"]))
    info_items.append(("电价(含税)", f"{tariff_usd:.5f} USD/kWh"))
    info_items.append(("电价(不含税)", f"{tariff_notax:.5f} USD/kWh"))
    info_items.append(("建设/运营", f"{inputs.basic.construction_months}个月 / {inputs.operational.operation_years}年"))
    info_items.append(("融资结构", f"资本金{inputs.financing.equity_ratio:.0%} | 利率{inputs.financing.long_term_loan_rate:.2%} | {inputs.financing.loan_term_years}年"))
    info_items.append(("所得税", f"{inputs.tax_financial.income_tax_rate:.0%}"))
    if pi.get("description"):
        info_items.append(("项目说明", pi["description"]))

    y = Y0 + Inches(0.4)
    for lb, vl in info_items:
        _txt(Inches(0.55), y, Inches(1.6), Inches(0.25), lb + ":", sz=9, c=GRAY, b=True)
        _txt(Inches(2.15), y, Inches(3.8), Inches(0.25), vl, sz=9, c=BLACK)
        y += Inches(0.25)

    # ── 左侧: 投资造价 ──
    y += Inches(0.15)
    _rect(Inches(0.3), y, Inches(0.12), Inches(0.32), fc=GREEN)
    _txt(Inches(0.55), y, Inches(3), Inches(0.32), "投资造价", sz=14, c=DARK, b=True)
    y += Inches(0.42)

    inv_items = [
        ("单位千瓦投资", f"{inputs.investment.resolve_unit_investment():,.0f} USD/kW"),
        ("静态总投资", f"{inputs.total_static_investment / 1e6:,.1f} M USD"),
        ("动态总投资", f"{inputs.total_dynamic_investment / 1e6:,.1f} M USD"),
        ("项目总投资", f"{inputs.total_investment / 1e6:,.1f} M USD"),
    ]
    for lb, vl in inv_items:
        _txt(Inches(0.55), y, Inches(1.6), Inches(0.25), lb + ":", sz=9, c=GRAY, b=True)
        _txt(Inches(2.15), y, Inches(3.8), Inches(0.25), vl, sz=9, c=BLACK)
        y += Inches(0.25)

    # ── 右侧: 财务指标 ──
    RX = Inches(6.3)
    _rect(RX - Inches(0.1), Y0, Inches(0.12), Inches(0.32), fc=ORANGE)
    _txt(RX + Inches(0.15), Y0, Inches(5), Inches(0.32),
         "财务指标", sz=14, c=DARK, b=True)

    irr_pre = result.project_irr_before_tax * 100
    irr_post = result.project_irr_after_tax * 100
    eq_irr = result.equity_irr * 100
    lcoe_val = result.lcoe
    pb = result.payback_after_tax
    npv_val = result.project_npv_after_tax / 1e6

    fin_rows = [
        ("全投资IRR(税前)", f"{irr_pre:.2f}%", LBLUE, True),
        ("全投资IRR(税后)", f"{irr_post:.2f}%", LBLUE, True),
        ("资本金IRR", f"{eq_irr:.2f}%", LGREEN, True),
        ("LCOE (USD/kWh)", f"{lcoe_val:.5f}", None, False),
        ("LCOE (元/kWh)", f"{lcoe_val * 7.1:.4f}", None, False),
        ("税后回收期 (年)", f"{pb:.2f}", None, False),
        ("NPV 税后 (M USD)", f"{npv_val:.1f}", LGREEN, True),
    ]

    y_t = Y0 + Inches(0.45)
    tbl = slide.shapes.add_table(
        len(fin_rows) + 1, 2, RX - Inches(0.1), y_t,
        Inches(6.9), Inches(len(fin_rows) * 0.35 + 0.32)).table
    tbl.columns[0].width = Inches(2.5)
    tbl.columns[1].width = Inches(4.4)

    _cell(tbl, 0, 0, "指标", fc=DARK, color=WHITE, bold=True, sz=9)
    _cell(tbl, 0, 1, "数值", fc=DARK, color=WHITE, bold=True, sz=9)

    for ri, (label, value, bg, bold) in enumerate(fin_rows):
        _cell(tbl, ri + 1, 0, label, sz=9, bold=True, al=PP_ALIGN.LEFT, fc=LGRAY)
        cl = RED_HL if bold else BLACK
        _cell(tbl, ri + 1, 1, value, sz=9, bold=bold, color=cl, fc=bg)

    # ── 底部: 数据来源 & 假设 ──
    y_bot = Inches(6.55)
    _rect(Inches(0), y_bot, SLIDE_W, Inches(0.95), fc=RGBColor(0xF0, 0xF4, 0xF8))
    _txt(Inches(0.5), y_bot + Inches(0.08), Inches(12), Inches(0.22),
         "数据来源 & 假设", sz=8, c=GRAY, b=True)
    _txt(Inches(0.5), y_bot + Inches(0.3), Inches(12), Inches(0.22),
         f"电价(含税): {tariff_usd:.5f} USD/kWh | "
         f"P90发电量: {inputs.basic.full_load_hours}h | "
         f"线损率: {inputs.basic.loss_rate:.2%} | "
         f"年上网电量: {inputs.net_annual_generation_mwh:,.0f} MWh",
         sz=8, c=BLACK)
    _txt(Inches(0.5), y_bot + Inches(0.52), Inches(12), Inches(0.22),
         f"折旧{inputs.operational.depreciation_years}年 | "
         f"残值率{inputs.operational.residual_rate:.1%} | "
         f"所得税{inputs.tax_financial.income_tax_rate:.0%} | "
         f"增值税{inputs.tax_financial.vat_rate:.0%} | "
         f"折现率{inputs.tax_financial.discount_rate:.1%}",
         sz=8, c=GRAY)
    _txt(Inches(0.5), y_bot + Inches(0.72), Inches(12), Inches(0.22),
         f"资本金{inputs.financing.equity_ratio:.0%} | "
         f"贷款利率{inputs.financing.long_term_loan_rate:.2%} | "
         f"贷款期限{inputs.financing.loan_term_years}年 | "
         f"运营期{inputs.operational.operation_years}年",
         sz=8, c=GRAY)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_comparison_ppt_bytes(
    projects_list: list,
    project_info: Optional[dict] = None,
) -> bytes:
    """生成多项目对比 PPT (支持 2-8 个方案)，返回二进制内容。
    projects_list: [(inputs, result), ...] 或 [(inputs, result, name_override), ...]
    project_info 可选，含 location/tariff_source/description 等自定义项目概况。"""
    from io import BytesIO
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK   = RGBColor(0x33, 0x33, 0x33)
    DARK    = RGBColor(0x1B, 0x2A, 0x4A)
    BLUE    = RGBColor(0x00, 0x6E, 0xB8)
    LBLUE   = RGBColor(0xD6, 0xEA, 0xF8)
    GREEN   = RGBColor(0x27, 0xAE, 0x60)
    LGREEN  = RGBColor(0xE8, 0xF8, 0xE8)
    GRAY    = RGBColor(0x7F, 0x8C, 0x8D)
    LGRAY   = RGBColor(0xF5, 0xF5, 0xF5)
    ORANGE  = RGBColor(0xE6, 0x7E, 0x22)
    RED_HL  = RGBColor(0xE7, 0x4C, 0x3C)

    N = len(projects_list)
    all_inputs = []
    all_results = []
    all_names = []
    for item in projects_list:
        all_inputs.append(item[0])
        all_results.append(item[1])
        all_names.append(item[2] if len(item) > 2 else item[0].basic.project_name)

    fs_data = max(7, 10 - N)  # adaptive font size for data cells
    fs_hdr = max(7, min(9, 11 - N))

    SLIDE_W_IN = 13.333
    SLIDE_W = Inches(SLIDE_W_IN)
    SLIDE_H = Inches(7.5)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE

    def _rect(l, t, w, h, fc=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
        s.line.fill.background()
        if fc:
            s.fill.solid()
            s.fill.fore_color.rgb = fc
        return s

    def _txt(l, t, w, h, text, sz=12, c=BLACK, b=False, al=PP_ALIGN.LEFT):
        tx = slide.shapes.add_textbox(l, t, w, h)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = "Microsoft YaHei"
        p.alignment = al

    def _cell(tbl, r, c, text, sz=9, color=BLACK, bold=False,
              al=PP_ALIGN.CENTER, fc=None):
        cl = tbl.cell(r, c)
        cl.text = ""
        p = cl.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = "Microsoft YaHei"
        p.alignment = al
        cl.vertical_anchor = MSO_ANCHOR.MIDDLE
        if fc:
            cl.fill.solid()
            cl.fill.fore_color.rgb = fc

    title_str = " vs ".join(all_names[:3])
    if N > 3:
        title_str = f"{all_names[0]} vs {all_names[1]} ... ({N}个方案)"

    # ── 标题栏 ──
    _rect(Inches(0), Inches(0), SLIDE_W, Inches(0.82), fc=DARK)
    _txt(Inches(0.4), Inches(0.05), Inches(9.5), Inches(0.4),
         title_str, sz=min(22, max(16, 24 - N)), c=WHITE, b=True)
    types_str = " / ".join(set(
        "海上" if inp.basic.project_type == "offshore" else "陆上"
        for inp in all_inputs))
    _txt(Inches(0.4), Inches(0.42), Inches(10), Inches(0.3),
         f"{types_str} · {N}方案经济性对比 · P90发电量 · "
         f"{all_inputs[0].operational.operation_years}年运营",
         sz=10, c=RGBColor(0xAA, 0xCC, 0xEE))

    # ── Logo (确保不超出边界) ──
    logo_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "mingyang_logo.png")
    if os.path.exists(logo_path):
        logo_h = Inches(0.5)
        logo_w = Inches(1.6)
        logo_left = Inches(SLIDE_W_IN - 1.6 - 0.3)
        slide.shapes.add_picture(logo_path, logo_left, Inches(0.15), width=logo_w, height=logo_h)

    # ── 投资造价对比表 (全宽) ──
    MARGIN = 0.35
    TBL_W = SLIDE_W_IN - 2 * MARGIN
    LABEL_W = 2.0 if N <= 4 else 1.7
    DATA_COL_W = (TBL_W - LABEL_W) / N

    Y_INV = 0.92
    _rect(Inches(MARGIN), Inches(Y_INV), Inches(0.1), Inches(0.22), fc=BLUE)
    _txt(Inches(MARGIN + 0.2), Inches(Y_INV), Inches(3), Inches(0.22),
         "投资造价对比", sz=12, c=DARK, b=True)

    inv_labels = ["配置", "P90 (h)", "CAPEX ($/kW)", "总投资 (M$)"]
    inv_data = []
    for inp in all_inputs:
        cap = inp.capacity_mw
        inv_data.append([
            f"{inp.basic.num_turbines}x{inp.basic.turbine_capacity_mw}MW={cap:.0f}MW",
            f"{inp.basic.full_load_hours:,}",
            f"{inp.investment.resolve_unit_investment():,.0f}",
            f"{inp.total_investment / 1e6:,.1f}",
        ])

    n_inv_rows = len(inv_labels)
    y_inv_tbl = Y_INV + 0.26
    t1 = slide.shapes.add_table(
        n_inv_rows + 1, N + 1,
        Inches(MARGIN), Inches(y_inv_tbl),
        Inches(TBL_W), Inches(n_inv_rows * 0.28 + 0.3)).table
    t1.columns[0].width = Inches(LABEL_W)
    for ci in range(N):
        t1.columns[ci + 1].width = Inches(DATA_COL_W)

    _cell(t1, 0, 0, "参数", fc=DARK, color=WHITE, bold=True, sz=fs_hdr)
    for ci, nm in enumerate(all_names):
        short = nm if len(nm) <= 18 else nm[:16] + ".."
        _cell(t1, 0, ci + 1, short, fc=DARK, color=WHITE, bold=True, sz=fs_hdr)

    for ri, lb in enumerate(inv_labels):
        bg = LGRAY if ri % 2 == 0 else None
        is_last = ri == n_inv_rows - 1
        _cell(t1, ri + 1, 0, lb, sz=fs_data, bold=True, al=PP_ALIGN.LEFT, fc=bg)
        for ci in range(N):
            fc_ = LGREEN if is_last else bg
            _cell(t1, ri + 1, ci + 1, inv_data[ci][ri], sz=fs_data, fc=fc_, bold=is_last)

    # ── 财务指标对比表 (全宽) ──
    y_fin_sec = y_inv_tbl + n_inv_rows * 0.28 + 0.3 + 0.15
    _rect(Inches(MARGIN), Inches(y_fin_sec), Inches(0.1), Inches(0.22), fc=ORANGE)
    _txt(Inches(MARGIN + 0.2), Inches(y_fin_sec), Inches(4), Inches(0.22),
         "财务指标对比 (P90)", sz=12, c=DARK, b=True)

    fin_labels = [
        ("全投IRR 税前", True, True),
        ("全投IRR 税后", True, True),
        ("资本金IRR", True, True),
        ("LCOE ($/kWh)", False, False),
        ("LCOE (元/kWh)", False, False),
        ("回收期 税后(年)", False, False),
        ("NPV 税后(M$)", True, True),
    ]
    fin_data_vals = []
    fin_data_strs = []
    for inp, res in zip(all_inputs, all_results):
        vals = [
            res.project_irr_before_tax * 100,
            res.project_irr_after_tax * 100,
            res.equity_irr * 100,
            res.lcoe,
            res.lcoe * 7.1,
            res.payback_after_tax,
            res.project_npv_after_tax / 1e6,
        ]
        strs = [
            f"{vals[0]:.2f}%",
            f"{vals[1]:.2f}%",
            f"{vals[2]:.2f}%",
            f"{vals[3]:.5f}",
            f"{vals[4]:.4f}",
            f"{vals[5]:.2f}",
            f"{vals[6]:.1f}",
        ]
        fin_data_vals.append(vals)
        fin_data_strs.append(strs)

    best_map = {}
    for ri, (_, is_key, higher_better) in enumerate(fin_labels):
        if not is_key:
            continue
        row_vals = [fin_data_vals[ci][ri] for ci in range(N)]
        best_idx = row_vals.index(max(row_vals) if higher_better else min(row_vals))
        if row_vals.count(row_vals[best_idx]) == 1:
            best_map[ri] = best_idx

    n_fin_rows = len(fin_labels)
    y_fin_tbl = y_fin_sec + 0.26
    t2 = slide.shapes.add_table(
        n_fin_rows + 1, N + 1,
        Inches(MARGIN), Inches(y_fin_tbl),
        Inches(TBL_W), Inches(n_fin_rows * 0.3 + 0.3)).table
    t2.columns[0].width = Inches(LABEL_W)
    for ci in range(N):
        t2.columns[ci + 1].width = Inches(DATA_COL_W)

    _cell(t2, 0, 0, "指标", fc=DARK, color=WHITE, bold=True, sz=fs_hdr)
    for ci, nm in enumerate(all_names):
        short = nm if len(nm) <= 18 else nm[:16] + ".."
        _cell(t2, 0, ci + 1, short, fc=DARK, color=WHITE, bold=True, sz=fs_hdr)

    bg_map = {0: LBLUE, 1: LBLUE, 2: LGREEN, 6: LGREEN}
    for ri, (lb, is_key, _) in enumerate(fin_labels):
        bg = bg_map.get(ri, LGRAY if ri % 2 == 0 else None)
        _cell(t2, ri + 1, 0, lb, sz=fs_data, bold=True, al=PP_ALIGN.LEFT, fc=LGRAY)
        for ci in range(N):
            is_best = best_map.get(ri, -1) == ci
            cl = RED_HL if is_best and is_key else BLACK
            extra = " *" if is_best and is_key else ""
            _cell(t2, ri + 1, ci + 1, fin_data_strs[ci][ri] + extra,
                  sz=fs_data, bold=is_key, color=cl, fc=bg)

    # ── 关键结论 ──
    y_conc = y_fin_tbl + n_fin_rows * 0.3 + 0.3 + 0.12
    _rect(Inches(MARGIN), Inches(y_conc), Inches(0.1), Inches(0.22), fc=RED_HL)
    _txt(Inches(MARGIN + 0.2), Inches(y_conc), Inches(3), Inches(0.22),
         "关键结论", sz=12, c=DARK, b=True)

    irr_posts = [r.project_irr_after_tax * 100 for r in all_results]
    best_idx = irr_posts.index(max(irr_posts))
    best_nm = all_names[best_idx]
    best_r = all_results[best_idx]
    lcoes = [r.lcoe for r in all_results]

    findings = [
        f"* 标记为各指标最优方案 | 全投IRR(税后)范围: {min(irr_posts):.2f}% ~ {max(irr_posts):.2f}%",
        f"{best_nm} 综合最优: IRR税后 {best_r.project_irr_after_tax*100:.2f}%, "
        f"LCOE {best_r.lcoe:.5f}$/kWh, 资本金IRR {best_r.equity_irr*100:.2f}%",
        f"LCOE范围: {min(lcoes):.5f} ~ {max(lcoes):.5f} USD/kWh "
        f"(差异 {(max(lcoes)-min(lcoes))*1000:.2f} $/MWh)",
    ]
    yy = y_conc + 0.28
    for fl in findings:
        _txt(Inches(MARGIN + 0.15), Inches(yy), Inches(12), Inches(0.22),
             f"▸ {fl}", sz=9, c=BLACK)
        yy += 0.24

    # ── 底部: 数据来源 & 假设 ──
    pi = project_info or {}
    y_bot = max(yy + 0.15, 6.3)
    _rect(Inches(0), Inches(y_bot), SLIDE_W, Inches(7.5 - y_bot), fc=RGBColor(0xF0, 0xF4, 0xF8))

    footer_parts = []
    if pi.get("location"):
        footer_parts.append(f"地点: {pi['location']}")
    if pi.get("tariff_source"):
        footer_parts.append(f"电价来源: {pi['tariff_source']}")
    tariff_str = " / ".join(f"{inp.tax_financial.tariff_with_tax:.4f}" for inp in all_inputs)
    footer_parts.append(f"电价(含税): {tariff_str} USD/kWh")
    ref_inp = all_inputs[0]
    footer_parts.append(
        f"资本金{ref_inp.financing.equity_ratio:.0%} | "
        f"贷款{ref_inp.financing.long_term_loan_rate:.2%}/{ref_inp.financing.loan_term_years}yr | "
        f"所得税{ref_inp.tax_financial.income_tax_rate:.0%} | "
        f"折旧{ref_inp.operational.depreciation_years}yr | "
        f"运营{ref_inp.operational.operation_years}yr")

    _txt(Inches(0.4), Inches(y_bot + 0.06), Inches(12.5), Inches(0.2),
         "数据来源 & 假设", sz=7, c=GRAY, b=True)
    f_y = y_bot + 0.24
    for fp in footer_parts:
        _txt(Inches(0.4), Inches(f_y), Inches(12.5), Inches(0.18), fp, sz=7, c=GRAY)
        f_y += 0.16

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def render_full_assessment(inputs: WindFarmFinancialInputs, result: CalculationResult, key_prefix: str = "main"):
    """
    渲染完整的项目评估视图（KPI + 参数 + 图表 + 明细表）。
    key_prefix 用于防止多次调用时 widget key 冲突。
    """
    # KPI
    plot_kpi_cards(result)
    st.markdown("---")

    # 国家参考
    profile = get_country_profile(inputs.basic.country)
    if profile:
        with st.expander(f"🌏 {profile.country_name_cn} 国别参数参考", expanded=False):
            c1, c2, c3, c4 = st.columns(4)
            c1.metric("企业所得税", f"{profile.corporate_income_tax_rate:.0%}")
            c2.metric("增值税", f"{profile.vat_rate:.0%}")
            c3.metric("典型贷款利率", f"{profile.typical_loan_rate:.2%}")
            c4.metric("典型资本金比例", f"{profile.typical_equity_ratio:.0%}")
            st.info(f"**优惠政策**: {profile.tax_incentive_description}")

    # 完整参数表
    with st.expander("📋 全部输入参数", expanded=False):
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            st.markdown("**基本信息**")
            st.write(f"- 项目名称: {inputs.basic.project_name}")
            st.write(f"- 类型: {'海上' if inputs.basic.project_type == 'offshore' else '陆上'}")
            st.write(f"- 国家: {inputs.basic.country}")
            st.write(f"- 装机: {inputs.capacity_mw:.0f} MW ({inputs.basic.num_turbines}×{inputs.basic.turbine_capacity_mw} MW)")
            st.write(f"- 满负荷小时数: {inputs.basic.full_load_hours} h")
            st.write(f"- 线损率: {inputs.basic.loss_rate:.2%}")
            st.write(f"- 年上网电量: {inputs.net_annual_generation_mwh:,.0f} MWh")
            st.write(f"- 建设期: {inputs.basic.construction_months} 月")
        with c2:
            st.markdown("**投资造价**")
            st.write(f"- 单位千瓦投资: {inputs.investment.resolve_unit_investment():,.1f} USD/kW")
            st.write(f"- 静态总投资: {inputs.total_static_investment / 1e6:,.1f} M$")
            st.write(f"- 建设期利息: {inputs.construction_interest / 1e6:,.2f} M$")
            st.write(f"- 动态总投资: {inputs.total_dynamic_investment / 1e6:,.1f} M$")
            st.write(f"- 流动资金: {inputs.working_capital / 1e6:,.2f} M$")
            st.write(f"- 项目总投资: {inputs.total_investment / 1e6:,.1f} M$")
        with c3:
            st.markdown("**融资条件**")
            st.write(f"- 资本金: {inputs.equity_amount / 1e6:,.1f} M$ ({inputs.financing.equity_ratio:.0%})")
            st.write(f"- 贷款: {inputs.debt_amount / 1e6:,.1f} M$")
            st.write(f"- 利率: {inputs.financing.long_term_loan_rate:.2%}")
            st.write(f"- 贷款期限: {inputs.financing.loan_term_years} 年")
            st.write(f"- 电价(含税): {inputs.tax_financial.tariff_with_tax:.4f} USD/kWh")
            st.write(f"- 电价(不含税): {inputs.tax_financial.tariff_without_tax:.4f} USD/kWh")
        with c4:
            st.markdown("**运营与税费**")
            st.write(f"- 运营期: {inputs.operational.operation_years} 年")
            st.write(f"- 折旧年限: {inputs.operational.depreciation_years} 年")
            st.write(f"- 残值率: {inputs.operational.residual_rate:.1%}")
            st.write(f"- 质保期: {inputs.operational.warranty.warranty_years} 年")
            st.write(f"- 保险费率: {inputs.operational.insurance_rate:.3%}")
            st.write(f"- 人员: {inputs.operational.staff_count} 人")
            st.write(f"- 增值税率: {inputs.tax_financial.vat_rate:.0%}")
            st.write(f"- 所得税率: {inputs.tax_financial.income_tax_rate:.0%}")
            st.write(f"- 折现率: {inputs.tax_financial.discount_rate:.1%}")

    # 投资造价分项明细
    with st.expander("💰 投资造价明细", expanded=False):
        if inputs.investment.onshore_detail:
            od = inputs.investment.onshore_detail
            inv_items = {
                "设备及安装工程": od.equipment_and_installation,
                "建筑工程": od.civil_works,
                "施工辅助工程": od.construction_auxiliary,
                "其他费用": od.other_costs,
                "储能工程": od.storage_cost,
                "送出线路": od.grid_connection_cost,
                "基本预备费": od.contingency,
                "**合计**": od.total_per_kw,
            }
            df_inv = pd.DataFrame([
                {"分项": k, "USD/kW": v, "总额(M$)": v * inputs.capacity_kw / 1e6}
                for k, v in inv_items.items() if v > 0 or k == "**合计**"
            ])
            st.dataframe(df_inv.style.format({"USD/kW": "{:,.1f}", "总额(M$)": "{:,.2f}"}),
                         use_container_width=True, hide_index=True)
        elif inputs.investment.offshore_detail:
            od = inputs.investment.offshore_detail
            inv_items = {
                "OEM (风机+塔筒)": od.oem_per_kw,
                "安装与运输": od.installation_per_kw,
                "基础工程": od.foundation_per_kw,
                "BOP 合计": od.bop.total_bop_per_kw,
                "**EPC 合计**": od.total_epc_per_kw,
            }
            df_inv = pd.DataFrame([
                {"分项": k, "USD/kW": v, "总额(M$)": v * inputs.capacity_kw / 1e6}
                for k, v in inv_items.items()
            ])
            st.dataframe(df_inv.style.format({"USD/kW": "{:,.1f}", "总额(M$)": "{:,.2f}"}),
                         use_container_width=True, hide_index=True)
        else:
            st.write(f"单位千瓦投资: {inputs.investment.resolve_unit_investment():,.1f} USD/kW")
            st.write(f"总投资: {inputs.total_investment / 1e6:,.1f} M$")

    # 运维成本逐年预览
    with st.expander("🔧 运维成本逐年预览", expanded=False):
        opex_rows = []
        sample_years = list(range(1, inputs.operational.operation_years + 1, max(1, inputs.operational.operation_years // 10)))
        if inputs.operational.operation_years not in sample_years:
            sample_years.append(inputs.operational.operation_years)
        for yr in sample_years:
            opex = inputs.operational.get_year_opex(yr, inputs.capacity_kw, inputs.total_static_investment)
            opex_rows.append({
                "运营年": yr,
                "材料费": opex["material"],
                "维修费": opex["repair"],
                "其他费用": opex["other"],
                "人员": opex["staff"],
                "保险": opex["insurance"],
                "海上专项": opex["offshore_extra"],
                "合计": sum(opex.values()),
            })
        df_opex = pd.DataFrame(opex_rows)
        st.dataframe(df_opex.style.format({c: "{:,.0f}" for c in df_opex.columns if c != "运营年"}),
                     use_container_width=True, hide_index=True)

    # 图表
    col_l, col_r = st.columns([2, 1])
    with col_l:
        plot_cashflow_chart(result)
    with col_r:
        plot_investment_breakdown(inputs)

    plot_profit_chart(result)

    st.markdown("---")
    plot_sensitivity(inputs)

    # 逐年明细
    st.markdown("---")
    st.markdown("### 📊 逐年现金流明细")
    df_data = []
    for f in result.annual_flows:
        df_data.append({
            "年份": f"建设-{f.year + 1}" if f.is_construction else f"运营-{f.year}",
            "营业收入": f.revenue,
            "补贴收入": f.subsidy_income,
            "经营成本": f.total_opex,
            "折旧": f.depreciation,
            "利息支出": f.loan_interest + f.wc_loan_interest,
            "总成本": f.total_cost,
            "附加税+资源税": f.surcharge + f.resource_tax,
            "利润总额": f.profit_before_tax,
            "所得税": f.income_tax,
            "净利润": f.net_profit,
            "全投资税后CF": f.project_net_cf_after_tax,
            "资本金CF": f.equity_net_cf,
        })
    df = pd.DataFrame(df_data)
    st.dataframe(
        df.style.format({c: "{:,.0f}" for c in df.columns if c != "年份"}),
        use_container_width=True, height=400,
    )

    # 项目概况(PPT用) + 导出按钮
    st.markdown("---")
    with st.expander("📝 PPT 项目概况 (选填，导出 PPT 时展示)", expanded=False):
        _pi_loc = st.text_input("项目地点", value="", placeholder="如：越南河静省 (Ha Tinh Province)", key=f"{key_prefix}_pi_loc")
        _pi_tariff = st.text_input("电价来源", value="", placeholder="如：Decision 1508/QĐ-BCT (2025.05.30)", key=f"{key_prefix}_pi_tariff")
        _pi_desc = st.text_area("项目说明/备注", value="", placeholder="如：近海风电(Nearshore)，2机型对比方案", height=68, key=f"{key_prefix}_pi_desc")
    _project_info = {}
    if _pi_loc:
        _project_info["location"] = _pi_loc
    if _pi_tariff:
        _project_info["tariff_source"] = _pi_tariff
    if _pi_desc:
        _project_info["description"] = _pi_desc

    dl_col1, dl_col2 = st.columns(2)
    with dl_col1:
        excel_bytes = export_to_excel(inputs, result)
        st.download_button(
            "📥 下载 Excel 报告", data=excel_bytes,
            file_name=f"{inputs.basic.project_name}_财务评价.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            key=f"{key_prefix}_dl_excel",
        )
    with dl_col2:
        ppt_bytes = generate_ppt_bytes(inputs, result, project_info=_project_info or None)
        st.download_button(
            "📥 下载 PPT 报告", data=ppt_bytes,
            file_name=f"{inputs.basic.project_name}_经济性评估.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"{key_prefix}_dl_ppt",
        )


def _rv_auto_params_box(inputs: WindFarmFinancialInputs, exclude: str = ""):
    """Display auto-extracted params from sidebar in a compact box"""
    b = inputs.basic
    inv = inputs.investment
    fin = inputs.financing
    tax = inputs.tax_financial
    ops = inputs.operational
    capex = inv.resolve_unit_investment()
    cap_mw = b.num_turbines * b.turbine_capacity_mw

    lines = []
    if exclude != "tariff":
        lines.append(f"Tariff: {tax.tariff_with_tax:.4f} USD/kWh")
    if exclude != "hours":
        lines.append(f"P90: {b.full_load_hours} h")
    lines.append(f"Capacity: {b.num_turbines} x {b.turbine_capacity_mw} MW = {cap_mw:.0f} MW")
    if exclude != "capex" and exclude != "turbine":
        lines.append(f"CAPEX: {capex:.0f} USD/kW")
    if exclude == "turbine":
        if inv.offshore_detail:
            non_t = capex - inv.offshore_detail.oem.turbine_price_per_kw
            lines.append(f"Non-turbine EPC: {non_t:.0f} USD/kW (auto)")
        elif inv.onshore_detail:
            non_eq = inv.onshore_detail.civil_works + inv.onshore_detail.other_costs + inv.onshore_detail.construction_auxiliary
            lines.append(f"BOP/Civil/Other: {non_eq:.0f} USD/kW (auto)")
        else:
            bop_est = capex * 0.40
            lines.append(f"BOP est.: {bop_est:.0f} USD/kW (40% of CAPEX)")
    lines.append(f"Equity: {fin.equity_ratio:.0%} | Rate: {fin.long_term_loan_rate:.1%} x {fin.loan_term_years}yr")
    lines.append(f"O&M: {ops.om_method} | Oper: {ops.operation_years}yr")
    lines.append(f"Tax: CIT {tax.income_tax_rate:.0%} | VAT {tax.vat_rate:.0%}")

    with st.expander("Auto-extracted from sidebar (no need to re-enter)", expanded=False):
        st.code("\n".join(lines), language=None)


def reverse_calc_panel(inputs: WindFarmFinancialInputs):
    is_offshore = inputs.basic.project_type == "offshore"
    has_offshore_epc = inputs.investment.offshore_detail is not None
    has_onshore_detail = inputs.investment.onshore_detail is not None

    st.info(
        "**How it works**: The solver automatically reads ALL parameters you entered in the sidebar. "
        "You only need to enter the **target value** below. The solver iterates to find the unknown."
    )

    tab1, tab2, tab3, tab4, tab5 = st.tabs([
        "IRR -> Tariff", "LCOE -> CAPEX", "IRR -> Hours",
        "LCOE -> Turbine Price", "IRR -> Turbine Price",
    ])
    irr_labels = {"project_before_tax": "Project Pre-tax", "project_after_tax": "Project After-tax", "equity": "Equity"}

    # ── Tab 1: IRR -> Tariff ──
    with tab1:
        st.markdown("#### Known: everything except tariff. Solve: tariff")
        st.caption("Scenario: You know the CAPEX, P90, financing, O&M etc. What tariff do you need to achieve target IRR?")
        _rv_auto_params_box(inputs, exclude="tariff")

        st.markdown("**Your input:**")
        c1, c2 = st.columns(2)
        t_irr = c1.number_input("Target IRR (%)", 1.0, 30.0, 8.0, step=0.5, key="rv1_irr") / 100.0
        t_type = c2.selectbox("IRR Type", list(irr_labels.keys()), 1, format_func=irr_labels.get, key="rv1_type")
        if st.button("Solve Tariff", key="rv1_btn", type="primary"):
            with st.spinner("Solving..."):
                try:
                    t = solve_tariff_for_target_irr(inputs, t_irr, t_type)
                    st.success(f"Required tariff (incl. tax): **{t:.5f} USD/kWh** ({t * 7.1:.4f} CNY/kWh)")
                except Exception as e:
                    st.error(f"Failed: {e}")

    # ── Tab 2: LCOE -> CAPEX ──
    with tab2:
        st.markdown("#### Known: tariff, P90, financing, O&M. Solve: max CAPEX")
        st.caption("Scenario: You know the tariff and all operating costs. What is the max investment you can afford?")
        _rv_auto_params_box(inputs, exclude="capex")

        st.markdown("**Your input:**")
        t_lcoe = st.number_input("Target LCOE (USD/kWh)", 0.001, 0.200, 0.030, step=0.001, format="%.4f", key="rv2_lcoe")
        if st.button("Solve CAPEX", key="rv2_btn", type="primary"):
            with st.spinner("Solving..."):
                try:
                    inv = solve_investment_for_target_lcoe(inputs, t_lcoe)
                    st.success(f"Max CAPEX: **{inv:.1f} USD/kW** ({inv * 7.1:.0f} CNY/kW)")
                except Exception as e:
                    st.error(f"Failed: {e}")

    # ── Tab 3: IRR -> Hours ──
    with tab3:
        st.markdown("#### Known: tariff, CAPEX, financing, O&M. Solve: min P90 hours")
        st.caption("Scenario: You know the investment and tariff. What is the minimum wind resource required?")
        _rv_auto_params_box(inputs, exclude="hours")

        st.markdown("**Your input:**")
        c1, c2 = st.columns(2)
        t_irr2 = c1.number_input("Target IRR (%)", 1.0, 30.0, 8.0, step=0.5, key="rv3_irr") / 100.0
        t_type2 = c2.selectbox("IRR Type", list(irr_labels.keys()), 1, format_func=irr_labels.get, key="rv3_type")
        if st.button("Solve Hours", key="rv3_btn", type="primary"):
            with st.spinner("Solving..."):
                try:
                    h = solve_hours_for_target_irr(inputs, t_irr2, t_type2)
                    st.success(f"Min full-load hours: **{h:.0f} h**")
                except Exception as e:
                    st.error(f"Failed: {e}")

    # ── Tab 4: LCOE -> Turbine Price ──
    with tab4:
        st.markdown("#### Solve: max turbine price (target = LCOE)")
        st.caption("You do NOT need to know or enter the turbine price - that is what this solver finds.")
        _rv_turbine_guide(is_offshore, has_offshore_epc, has_onshore_detail, inputs)
        _rv_auto_params_box(inputs, exclude="turbine")

        st.markdown("**Your input (only 1 field):**")
        t_lcoe4 = st.number_input(
            "Target LCOE (USD/kWh)", 0.001, 0.200, 0.040, step=0.001, format="%.4f", key="rv4_lcoe",
        )
        if st.button("Solve Turbine Price", key="rv4_btn", type="primary"):
            with st.spinner("Solving..."):
                try:
                    price = solve_turbine_price_for_target_lcoe(inputs, t_lcoe4)
                    if price is not None:
                        st.success(f"Max turbine OEM price: **{price:,.1f} USD/kW** ({price * 7.1:,.0f} CNY/kW)")
                        _rv_turbine_result_detail(is_offshore, has_offshore_epc, has_onshore_detail, inputs, price)
                    else:
                        st.error("Solve failed.")
                except Exception as e:
                    st.error(f"Failed: {e}")

    # ── Tab 5: IRR -> Turbine Price ──
    with tab5:
        st.markdown("#### Solve: max turbine price (target = IRR)")
        st.caption("You do NOT need to know or enter the turbine price - that is what this solver finds.")
        _rv_turbine_guide(is_offshore, has_offshore_epc, has_onshore_detail, inputs)
        _rv_auto_params_box(inputs, exclude="turbine")

        st.markdown("**Your input (only 2 fields):**")
        c1, c2 = st.columns(2)
        t_irr5 = c1.number_input("Target IRR (%)", 1.0, 30.0, 8.0, step=0.5, key="rv5_irr") / 100.0
        t_type5 = c2.selectbox("IRR Type", list(irr_labels.keys()), 1, format_func=irr_labels.get, key="rv5_type")
        if st.button("Solve Turbine Price", key="rv5_btn", type="primary"):
            with st.spinner("Solving..."):
                try:
                    price = solve_turbine_price_for_target_irr(inputs, t_irr5, t_type5)
                    if price is not None:
                        st.success(f"Max turbine OEM price: **{price:,.1f} USD/kW** ({price * 7.1:,.0f} CNY/kW)")
                        _rv_turbine_result_detail(is_offshore, has_offshore_epc, has_onshore_detail, inputs, price)
                    else:
                        st.error("Solve failed.")
                except Exception as e:
                    st.error(f"Failed: {e}")


def _rv_turbine_guide(is_offshore, has_offshore_epc, has_onshore_detail, inputs):
    """Show onshore vs offshore workflow guidance for turbine price solver."""
    inv = inputs.investment
    if has_offshore_epc:
        detail = inv.offshore_detail
        oem = detail.oem
        epc_total = detail.total_epc_per_kw
        non_turbine = epc_total - oem.turbine_price_per_kw
        tower_per_kw = detail.oem_per_kw - oem.turbine_price_per_kw
        st.info(
            f"**Current mode: Offshore (EPC breakdown)**\n\n"
            f"- EPC total: {epc_total:,.0f} USD/kW | Non-turbine: {non_turbine:,.0f} USD/kW\n"
            f"- Solver holds constant: tower (~{tower_per_kw:,.0f} USD/kW), "
            f"installation ({detail.installation_per_kw:,.0f}), "
            f"foundation ({detail.foundation_per_kw:,.0f}), "
            f"BOP ({detail.bop.total_bop_per_kw:,.0f})\n"
            f"- Solver adjusts: **turbine OEM price only**\n\n"
            f"Sidebar checklist: make sure tower, installation, foundation, BOP, cable "
            f"costs in EPC breakdown are correctly filled."
        )
    elif has_onshore_detail:
        on = inv.onshore_detail
        non_equip = on.civil_works + on.construction_auxiliary + on.other_costs
        contingency = (on.equipment_and_installation + on.civil_works + on.construction_auxiliary + on.other_costs) * on.contingency_rate
        total_fixed = non_equip + contingency + on.storage_cost + on.grid_connection_cost
        st.info(
            f"**Current mode: Onshore (investment breakdown)**\n\n"
            f"- Equipment & installation: {on.equipment_and_installation:,.0f} USD/kW "
            f"(contains turbine + non-turbine equipment)\n"
            f"- Fixed non-turbine: civil {on.civil_works:,.0f} + auxiliary {on.construction_auxiliary:,.0f} "
            f"+ other {on.other_costs:,.0f} + contingency ~{contingency:,.0f} "
            f"+ storage {on.storage_cost:,.0f} + grid {on.grid_connection_cost:,.0f} "
            f"= **{total_fixed:,.0f} USD/kW**\n"
            f"- Solver adjusts: **turbine portion within equipment**, "
            f"keeps civil/auxiliary/other/storage/grid constant\n\n"
            f"Sidebar checklist: make sure civil works, auxiliary, other costs, "
            f"storage, grid connection are correctly filled."
        )
    else:
        capex = inv.resolve_unit_investment()
        bop_est = capex * 0.40
        st.warning(
            f"**Current mode: Quick (no cost breakdown)**\n\n"
            f"- Total CAPEX: {capex:,.0f} USD/kW\n"
            f"- BOP estimated: {bop_est:,.0f} USD/kW (40% of CAPEX, held constant)\n"
            f"- Turbine estimated: {capex - bop_est:,.0f} USD/kW (60%, to be solved)\n\n"
            f"**For more accurate results**, switch to **Detailed mode** "
            f"in the sidebar and fill in the investment breakdown."
        )

    with st.expander("Onshore vs Offshore - what's different?", expanded=False):
        st.markdown("""
| | Offshore | Onshore | Quick (no breakdown) |
|---|---|---|---|
| **Cost structure** | EPC = OEM + tower + install + foundation + BOP + cable | Equipment & install + civil + auxiliary + other | Single CAPEX number |
| **What solver adjusts** | OEM turbine price only | Turbine portion within equipment | 60% of CAPEX |
| **What stays fixed** | Tower, install, foundation, BOP, cable, all other EPC items | Civil works, auxiliary, other costs, contingency, storage, grid | 40% of CAPEX (BOP estimate) |
| **Sidebar preparation** | Fill EPC breakdown (tower, install, foundation, BOP, cable costs) | Fill investment breakdown (civil, auxiliary, other, storage, grid) | Just total CAPEX |
| **Accuracy** | Highest (explicit cost split) | High (explicit non-turbine items) | Approximate (BOP is estimated) |
| **Recommended for** | Offshore wind farms with EPC contracts | Onshore wind farms with feasibility study data | Quick screening / early stage |
        """)


def _rv_turbine_result_detail(is_offshore, has_offshore_epc, has_onshore_detail, inputs, price):
    """Show detailed breakdown of how the solved turbine price fits into the cost structure."""
    inv = inputs.investment
    if has_offshore_epc:
        detail = inv.offshore_detail
        oem = detail.oem
        old_epc = detail.total_epc_per_kw
        new_epc = old_epc - oem.turbine_price_per_kw + price
        tower_per_kw = detail.oem_per_kw - oem.turbine_price_per_kw
        non_oem_epc = old_epc - detail.oem_per_kw
        st.info(
            f"**Cost breakdown with solved turbine price:**\n"
            f"- Turbine OEM: **{price:,.0f}** USD/kW (solved)\n"
            f"- Tower: ~{tower_per_kw:,.0f} USD/kW (unchanged)\n"
            f"- Install+Foundation+BOP: {non_oem_epc:,.0f} USD/kW (unchanged)\n"
            f"- New total EPC: **{new_epc:,.0f}** USD/kW"
        )
    elif has_onshore_detail:
        on = inv.onshore_detail
        if on.turbine_price_per_kw > 0:
            non_turbine_equip = on.non_turbine_equip_per_kw
        else:
            non_turbine_equip = on.equipment_and_installation * 0.30
        new_equip = price + non_turbine_equip
        non_equip = on.civil_works + on.construction_auxiliary + on.other_costs
        new_total = new_equip + non_equip
        new_total_with_cont = new_total * (1 + on.contingency_rate) + on.storage_cost + on.grid_connection_cost
        st.info(
            f"**Cost breakdown with solved turbine price:**\n"
            f"- Turbine: **{price:,.0f}** USD/kW (solved)\n"
            f"- Non-turbine equipment: {non_turbine_equip:,.0f} USD/kW (unchanged)\n"
            f"- New equipment & install: **{new_equip:,.0f}** USD/kW\n"
            f"- Civil + auxiliary + other: {non_equip:,.0f} USD/kW (unchanged)\n"
            f"- Estimated new total: ~**{new_total_with_cont:,.0f}** USD/kW"
        )
    else:
        capex = inv.resolve_unit_investment()
        bop = capex * 0.40
        new_capex = price + bop
        st.info(
            f"**Estimated breakdown:**\n"
            f"- Turbine: **{price:,.0f}** USD/kW (solved)\n"
            f"- BOP (estimated 40%): {bop:,.0f} USD/kW (unchanged)\n"
            f"- Estimated new CAPEX: **{new_capex:,.0f}** USD/kW"
        )


# ════════════════════════════════════════════════════════════════════════════
# 各国市场概览面板
# ════════════════════════════════════════════════════════════════════════════

def render_market_overview():
    """各国风电市场概览 — IRR / WACC 对比 + 国家详情 + 汇总表"""
    st.header("🌍 各国风电市场概览")
    st.caption("数据来源: BNEF, IRENA, IEA, World Bank, 各国政府部门 | 仅供参考")

    profiles = list(_PROFILES.values())

    # ─── 顶部: 各国 Equity IRR 横向对比柱状图 ───
    st.subheader("📊 各国 Equity IRR 对比")
    _render_irr_comparison_chart(profiles)

    st.markdown("---")

    # ─── 中间: 国家选择器 + 详情卡片 ───
    st.subheader("🔍 国家详情")
    country_options = {p.country_name_cn: p for p in profiles}
    selected_name = st.selectbox(
        "选择国家", list(country_options.keys()), key="market_country_select"
    )
    if selected_name:
        _render_country_detail_card(country_options[selected_name])

    st.markdown("---")

    # ─── 电价对比表 ───
    st.subheader("⚡ 各国电价参考对比")
    _tariff_rows = get_all_tariff_summary()
    if _tariff_rows:
        _tdf = pd.DataFrame(_tariff_rows)
        _tdf_display = _tdf.rename(columns={
            "country_cn": "国家", "onshore_low": "陆上下限", "onshore_high": "陆上上限",
            "offshore_low": "海上下限", "offshore_high": "海上上限",
            "onshore_source": "陆上来源", "offshore_source": "海上来源", "mechanism": "定价机制",
        })[["国家", "陆上下限", "陆上上限", "海上下限", "海上上限", "定价机制", "陆上来源"]]
        st.dataframe(_tdf_display, use_container_width=True, hide_index=True)
        st.caption("单位: USD/kWh | 数据自动更新（IRENA优先，静态数据兜底）")
    else:
        st.info("暂无电价数据")

    st.markdown("---")

    # ─── 底部: 全部国家对比汇总表 ───
    st.subheader("📋 全部国家对比汇总表")
    _render_summary_table(profiles)


def _render_irr_comparison_chart(profiles: list):
    """各国 Equity IRR 范围对比 — Plotly grouped bar chart"""
    rows = []
    for p in profiles:
        onshore_irrs = [b for b in p.benchmarks if b.metric == "Equity IRR" and b.project_type in ("onshore", "all")]
        offshore_irrs = [b for b in p.benchmarks if b.metric == "Equity IRR" and b.project_type in ("offshore", "nearshore", "all")]

        if onshore_irrs:
            low = min(b.value_low for b in onshore_irrs)
            high = max(b.value_high for b in onshore_irrs)
            rows.append({"国家": p.country_name_cn, "类型": "陆上", "IRR下限(%)": low, "IRR上限(%)": high, "中值(%)": (low + high) / 2})
        if offshore_irrs:
            low = min(b.value_low for b in offshore_irrs)
            high = max(b.value_high for b in offshore_irrs)
            rows.append({"国家": p.country_name_cn, "类型": "海上/近海", "IRR下限(%)": low, "IRR上限(%)": high, "中值(%)": (low + high) / 2})

    if not rows:
        st.info("暂无 IRR 基准数据")
        return

    df = pd.DataFrame(rows)

    fig = go.Figure()
    colors = {"陆上": "#2ecc71", "海上/近海": "#3498db"}
    for typ in df["类型"].unique():
        sub = df[df["类型"] == typ]
        fig.add_trace(go.Bar(
            name=typ,
            x=sub["国家"],
            y=sub["中值(%)"],
            error_y=dict(
                type="data",
                symmetric=False,
                array=(sub["IRR上限(%)"] - sub["中值(%)"]).tolist(),
                arrayminus=(sub["中值(%)"] - sub["IRR下限(%)"]).tolist(),
            ),
            marker_color=colors.get(typ, "#95a5a6"),
            text=sub.apply(lambda r: f"{r['IRR下限(%)']:.0f}-{r['IRR上限(%)']:.0f}%", axis=1),
            textposition="outside",
        ))

    fig.update_layout(
        barmode="group",
        yaxis_title="Equity IRR (%)",
        xaxis_title="",
        height=450,
        legend=dict(orientation="h", yanchor="bottom", y=1.02, xanchor="right", x=1),
        margin=dict(t=60, b=40),
    )
    st.plotly_chart(fig, use_container_width=True)


def _source_link(name: str, url: str) -> str:
    """生成来源 markdown 链接，无 URL 则返回纯文本"""
    if url:
        return f"[{name}]({url})"
    return name


_TYPE_CN = {"onshore": "陆上", "offshore": "海上", "nearshore": "近海", "all": "通用"}


def _render_country_detail_card(p: CountryProfile):
    """单个国家的详情卡片 — 分板块展示"""
    st.markdown(f"### {p.country_name_cn} ({p.country_name})")

    # ─── 基本参数指标卡 ───
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("货币", p.currency)
    col2.metric("资本金比例", f"{p.typical_equity_ratio*100:.0f}%")
    col3.metric("贷款利率", f"{p.typical_loan_rate*100:.1f}%")
    col4.metric("企业所得税", f"{p.corporate_income_tax_rate*100:.1f}%")

    col5, col6, col7, col8 = st.columns(4)
    col5.metric("增值税", f"{p.vat_rate*100:.0f}%")
    col6.metric("贷款期限", f"{p.typical_loan_term} 年")
    if p.onshore_tariff_range != (0.0, 0.0):
        col7.metric("陆上电价", f"{p.onshore_tariff_range[0]*100:.1f}-{p.onshore_tariff_range[1]*100:.1f} ¢/kWh")
    else:
        col7.metric("陆上电价", "—")
    if p.offshore_tariff_range != (0.0, 0.0):
        col8.metric("海上电价", f"{p.offshore_tariff_range[0]*100:.1f}-{p.offshore_tariff_range[1]*100:.1f} ¢/kWh")
    else:
        col8.metric("海上电价", "—")

    if p.has_wind_tax_incentive:
        st.info(f"💡 **税收优惠**: {p.tax_incentive_description}")

    # ─── 分板块市场报告 ───
    rpt = getattr(p, "market_report", None)
    has_report = rpt and (rpt.official_benchmarks or rpt.bnef_hurdles
                          or rpt.actual_cases or rpt.wacc_data or rpt.summary)

    if has_report:
        _render_market_report_sections(rpt)

        _render_radar_chart(p)

    elif p.benchmarks:
        st.markdown("#### 📈 市场基准数据")
        bm_rows = []
        for b in p.benchmarks:
            range_str = (f"{b.value_low:.1f}%" if b.value_low == b.value_high
                         else f"{b.value_low:.1f}% - {b.value_high:.1f}%")
            source_str = _source_link(b.source, b.source_url)
            bm_rows.append({
                "指标": b.metric,
                "项目类型": _TYPE_CN.get(b.project_type, b.project_type),
                "数值范围": range_str,
                "来源": source_str,
                "年份": b.year,
                "备注": b.note,
            })
        md_header = "| 指标 | 项目类型 | 数值范围 | 来源 | 年份 | 备注 |\n| :-- | :-- | :-- | :-- | :-- | :-- |"
        md_body = "\n".join(
            f"| {r['指标']} | {r['项目类型']} | {r['数值范围']} | {r['来源']} | {r['年份']} | {r['备注']} |"
            for r in bm_rows
        )
        st.markdown(md_header + "\n" + md_body, unsafe_allow_html=True)

        _render_radar_chart(p)
    else:
        st.warning("该国家暂无市场基准数据")

    st.caption(f"数据更新: {p.data_updated}")


def _render_market_report_sections(rpt):
    """渲染分板块市场报告（官方基准 / BNEF / 实际案例 / WACC / 总结）"""

    # ─── 一、官方 IRR 基准 ───
    if rpt.official_benchmarks:
        st.markdown("#### 一、官方/政府定价模型中的 IRR 基准")
        header = "| 来源 | 项目类型 | 股权收益率 (Equity IRR) | 融资结构 | 备注 |\n| :-- | :-- | :-- | :-- | :-- |"
        rows = []
        sources = []
        for b in rpt.official_benchmarks:
            rows.append(
                f"| {b.source_ref} | {_TYPE_CN.get(b.project_type, b.project_type)} "
                f"| {b.equity_irr:.1f}% | {b.financing_structure} | {b.notes} |"
            )
            sources.append(_source_link(b.source, b.source_url))
        st.markdown(header + "\n" + "\n".join(rows), unsafe_allow_html=True)
        unique_sources = list(dict.fromkeys(sources))
        st.caption("来源: " + " ; ".join(unique_sources))

    # ─── 二、BNEF Hurdle IRR ───
    if rpt.bnef_hurdles:
        st.markdown("#### 二、国际研报 Hurdle IRR")
        header = "| 项目类型 | 当前值 | 2030年(预测) | 2050年(预测) | 融资假设 |\n| :-- | :-- | :-- | :-- | :-- |"
        rows = []
        sources = []
        for b in rpt.bnef_hurdles:
            rows.append(
                f"| {_TYPE_CN.get(b.project_type, b.project_type)} "
                f"| {b.current_value} | {b.forecast_2030 or '—'} "
                f"| {b.forecast_2050 or '—'} | {b.financing_assumption or '—'} |"
            )
            sources.append(_source_link(f"{b.source} ({b.report_name}, {b.year})", b.source_url))
        st.markdown(header + "\n" + "\n".join(rows), unsafe_allow_html=True)
        unique_sources = list(dict.fromkeys(sources))
        st.caption("来源: " + " ; ".join(unique_sources))

    # ─── 三、实际项目案例 ───
    if rpt.actual_cases:
        st.markdown("#### 三、实际项目 IRR 案例")
        header = "| 项目/公司 | 类型 | IRR 类型 | 实际/预期 IRR | 备注 |\n| :-- | :-- | :-- | :-- | :-- |"
        rows = []
        sources = []
        for c in rpt.actual_cases:
            rows.append(
                f"| {c.project_name} | {_TYPE_CN.get(c.project_type, c.project_type)} "
                f"| {c.irr_type} | {c.irr_value} | {c.notes} |"
            )
            if c.source:
                sources.append(_source_link(c.source, c.source_url))
        st.markdown(header + "\n" + "\n".join(rows), unsafe_allow_html=True)
        if sources:
            unique_sources = list(dict.fromkeys(sources))
            st.caption("来源: " + " ; ".join(unique_sources))

    # ─── 四、WACC / 融资成本 ───
    if rpt.wacc_data:
        st.markdown("#### 四、WACC / 融资成本")
        header = "| 来源 | 指标 | 数值 | 备注 |\n| :-- | :-- | :-- | :-- |"
        rows = []
        sources = []
        for w in rpt.wacc_data:
            rows.append(f"| {w.source} | {w.indicator} | {w.value} | {w.notes} |")
            sources.append(_source_link(w.source, w.source_url))
        st.markdown(header + "\n" + "\n".join(rows), unsafe_allow_html=True)
        unique_sources = list(dict.fromkeys(sources))
        st.caption("来源: " + " ; ".join(unique_sources))

    # ─── 五、总结判断 ───
    if rpt.summary:
        st.markdown("#### 五、总结判断")
        header = "| 场景 | IRR 范围 |\n| :-- | :-- |"
        rows = [f"| {s.scenario} | {s.irr_range} |" for s in rpt.summary]
        st.markdown(header + "\n" + "\n".join(rows), unsafe_allow_html=True)
        if rpt.summary_conclusion:
            st.markdown(f"> {rpt.summary_conclusion}")


def _render_radar_chart(p: CountryProfile):
    """国家维度雷达图: IRR / WACC / 电价 / 税率 / 贷款利率"""
    equity_irrs = [b for b in p.benchmarks if b.metric == "Equity IRR"]
    wacc_list = [b for b in p.benchmarks if b.metric == "WACC"]

    avg_irr = np.mean([(b.value_low + b.value_high) / 2 for b in equity_irrs]) if equity_irrs else 0
    avg_wacc = np.mean([(b.value_low + b.value_high) / 2 for b in wacc_list]) if wacc_list else 0

    on_tariff_mid = (p.onshore_tariff_range[0] + p.onshore_tariff_range[1]) / 2 * 100 if p.onshore_tariff_range != (0.0, 0.0) else 0
    off_tariff_mid = (p.offshore_tariff_range[0] + p.offshore_tariff_range[1]) / 2 * 100 if p.offshore_tariff_range != (0.0, 0.0) else 0
    tariff_mid = max(on_tariff_mid, off_tariff_mid)

    categories = ["Equity IRR(%)", "WACC(%)", "电价(¢/kWh)", "所得税(%)", "贷款利率(%)"]
    values = [avg_irr, avg_wacc, tariff_mid, p.corporate_income_tax_rate * 100, p.typical_loan_rate * 100]
    values.append(values[0])
    categories.append(categories[0])

    fig = go.Figure(data=go.Scatterpolar(
        r=values,
        theta=categories,
        fill="toself",
        name=p.country_name_cn,
        line_color="#3498db",
    ))
    fig.update_layout(
        polar=dict(radialaxis=dict(visible=True, range=[0, max(values) * 1.2 if max(values) > 0 else 20])),
        showlegend=False,
        height=380,
        margin=dict(t=30, b=30),
    )
    st.plotly_chart(fig, use_container_width=True)


def _render_summary_table(profiles: list):
    """全部国家对比汇总表"""
    rows = []
    for p in profiles:
        equity_irrs = [b for b in p.benchmarks if b.metric == "Equity IRR"]
        wacc_list = [b for b in p.benchmarks if b.metric == "WACC"]
        project_irrs = [b for b in p.benchmarks if b.metric == "Project IRR"]

        def _range_str(blist):
            if not blist:
                return "—"
            low = min(b.value_low for b in blist)
            high = max(b.value_high for b in blist)
            return f"{low:.1f}-{high:.1f}%"

        on_tariff = f"{p.onshore_tariff_range[0]*100:.1f}-{p.onshore_tariff_range[1]*100:.1f}" if p.onshore_tariff_range != (0.0, 0.0) else "—"
        off_tariff = f"{p.offshore_tariff_range[0]*100:.1f}-{p.offshore_tariff_range[1]*100:.1f}" if p.offshore_tariff_range != (0.0, 0.0) else "—"

        rows.append({
            "国家": p.country_name_cn,
            "Equity IRR": _range_str(equity_irrs),
            "Project IRR": _range_str(project_irrs),
            "WACC": _range_str(wacc_list),
            "陆上电价(¢/kWh)": on_tariff,
            "海上电价(¢/kWh)": off_tariff,
            "所得税": f"{p.corporate_income_tax_rate*100:.1f}%",
            "贷款利率": f"{p.typical_loan_rate*100:.1f}%",
            "资本金比例": f"{p.typical_equity_ratio*100:.0f}%",
            "税收优惠": "✅" if p.has_wind_tax_incentive else "❌",
        })

    df = pd.DataFrame(rows)
    st.dataframe(df, use_container_width=True, hide_index=True, height=450)


# ════════════════════════════════════════════════════════════════════════════
# 项目对比面板
# ════════════════════════════════════════════════════════════════════════════

def comparison_page():
    """多项目对比页面"""
    st.header("📊 项目对比")

    projects = st.session_state.projects
    if len(projects) < 2:
        st.warning("至少需要保存 **2 个项目** 才能进行对比。请在「项目评估」页面保存项目后再来。")
        return

    options = {pid: f"{p['name']} ({p['saved_at']})" for pid, p in projects.items()}
    selected = st.multiselect(
        "选择要对比的项目 (至少选 2 个)",
        options=list(options.keys()),
        format_func=lambda x: options[x],
        default=list(options.keys())[:min(4, len(options))],
    )

    if len(selected) < 2:
        st.info("请至少选择 2 个项目。")
        return

    # ──── 对比 PPT 下载 (2~8 个方案) ────
    with st.expander("📝 对比 PPT 项目概况 (选填)", expanded=False):
        _cpi_loc = st.text_input("项目地点", value="", placeholder="如：越南河静省", key="cpi_loc")
        _cpi_tariff = st.text_input("电价来源", value="", placeholder="如：Decision 1508", key="cpi_tariff")
        _cpi_desc = st.text_area("项目说明", value="", placeholder="如：近海风电 多方案对比", height=68, key="cpi_desc")
    _cpi = {}
    if _cpi_loc:
        _cpi["location"] = _cpi_loc
    if _cpi_tariff:
        _cpi["tariff_source"] = _cpi_tariff
    if _cpi_desc:
        _cpi["description"] = _cpi_desc

    ppt_projects = [
        (projects[pid]["inputs"], projects[pid]["result"], projects[pid]["name"])
        for pid in selected
    ]
    ppt_bytes = generate_comparison_ppt_bytes(ppt_projects, project_info=_cpi or None)
    fname = "_vs_".join(projects[pid]["name"] for pid in selected[:3])
    if len(selected) > 3:
        fname += f"_等{len(selected)}方案"
    fname += "_对比.pptx"
    st.download_button(
        f"📥 下载 {len(selected)} 方案对比 PPT", data=ppt_bytes,
        file_name=fname,
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        key="dl_comparison_ppt",
    )

    # ──── KPI 对比表 ────
    st.markdown("### 关键指标对比")
    rows = []
    for pid in selected:
        p = projects[pid]
        r: CalculationResult = p["result"]
        inp: WindFarmFinancialInputs = p["inputs"]
        rows.append({
            "项目": p["name"],
            "类型": "海上" if inp.basic.project_type == "offshore" else "陆上",
            "容量(MW)": inp.capacity_mw,
            "小时数(h)": inp.basic.full_load_hours,
            "投资(USD/kW)": inp.investment.resolve_unit_investment(),
            "电价(USD/kWh)": inp.tax_financial.tariff_with_tax,
            "全投资IRR(税后)": r.project_irr_after_tax,
            "资本金IRR": r.equity_irr,
            "LCOE(USD/kWh)": r.lcoe,
            "回收期(年)": r.payback_after_tax,
            "NPV(M$)": r.project_npv_after_tax / 1e6,
            "ROI": r.roi,
            "ROE": r.roe,
        })

    df = pd.DataFrame(rows)
    fmt_map = {
        "投资(USD/kW)": "{:,.1f}",
        "电价(USD/kWh)": "{:.4f}",
        "全投资IRR(税后)": "{:.2%}",
        "资本金IRR": "{:.2%}",
        "LCOE(USD/kWh)": "{:.5f}",
        "回收期(年)": "{:.1f}",
        "NPV(M$)": "{:,.1f}",
        "ROI": "{:.2%}",
        "ROE": "{:.2%}",
    }
    st.dataframe(df.style.format(fmt_map), use_container_width=True, hide_index=True)

    # ──── 柱状图对比 ────
    st.markdown("### 核心指标并列对比")

    col1, col2 = st.columns(2)

    with col1:
        fig = go.Figure()
        names = [projects[pid]["name"] for pid in selected]
        irr_vals = [projects[pid]["result"].project_irr_after_tax * 100 for pid in selected]
        eq_irr_vals = [projects[pid]["result"].equity_irr * 100 for pid in selected]
        fig.add_trace(go.Bar(x=names, y=irr_vals, name="全投资IRR(税后)", marker_color="#2E75B6"))
        fig.add_trace(go.Bar(x=names, y=eq_irr_vals, name="资本金IRR", marker_color="#548235"))
        fig.update_layout(title="IRR 对比 (%)", barmode="group", height=380, yaxis_title="%")
        st.plotly_chart(fig, use_container_width=True)

    with col2:
        fig = go.Figure()
        lcoe_vals = [projects[pid]["result"].lcoe * 1000 for pid in selected]
        fig.add_trace(go.Bar(x=names, y=lcoe_vals, name="LCOE", marker_color="#BF8F00"))
        fig.update_layout(title="LCOE 对比 (USD/MWh)", height=380, yaxis_title="USD/MWh")
        st.plotly_chart(fig, use_container_width=True)

    # ──── 雷达图 ────
    st.markdown("### 综合能力雷达图")

    categories = ["IRR", "ROI", "ROE", "回收期(短优)", "LCOE(低优)"]
    fig = go.Figure()

    all_irr = [projects[pid]["result"].project_irr_after_tax for pid in selected]
    all_roi = [projects[pid]["result"].roi for pid in selected]
    all_roe = [projects[pid]["result"].roe for pid in selected]
    all_payback = [projects[pid]["result"].payback_after_tax for pid in selected]
    all_lcoe = [projects[pid]["result"].lcoe for pid in selected]

    def normalize(vals, invert=False):
        mn, mx = min(vals), max(vals)
        if mx == mn:
            return [0.5] * len(vals)
        if invert:
            return [(mx - v) / (mx - mn) for v in vals]
        return [(v - mn) / (mx - mn) for v in vals]

    n_irr = normalize(all_irr)
    n_roi = normalize(all_roi)
    n_roe = normalize(all_roe)
    n_pay = normalize(all_payback, invert=True)
    n_lcoe = normalize(all_lcoe, invert=True)

    for i, pid in enumerate(selected):
        vals = [n_irr[i], n_roi[i], n_roe[i], n_pay[i], n_lcoe[i]]
        vals.append(vals[0])
        fig.add_trace(go.Scatterpolar(
            r=vals,
            theta=categories + [categories[0]],
            name=projects[pid]["name"],
            fill="toself",
            opacity=0.6,
            line_color=COLOR_PALETTE[i % len(COLOR_PALETTE)],
        ))

    fig.update_layout(
        polar=dict(radialaxis=dict(visible=True, range=[0, 1])),
        title="项目综合能力雷达图 (归一化)",
        height=500,
    )
    st.plotly_chart(fig, use_container_width=True)

    # ──── 现金流对比 ────
    st.markdown("### 累计现金流对比")
    fig = go.Figure()
    for i, pid in enumerate(selected):
        p = projects[pid]
        cum, running = [], 0.0
        for f in p["result"].annual_flows:
            running += f.project_net_cf_after_tax
            cum.append(running)
        fig.add_trace(go.Scatter(
            y=cum, name=p["name"],
            line=dict(color=COLOR_PALETTE[i % len(COLOR_PALETTE)], width=2.5),
        ))
    fig.update_layout(title="累计净现金流", xaxis_title="年份", yaxis_title="USD", height=400)
    st.plotly_chart(fig, use_container_width=True)


# ════════════════════════════════════════════════════════════════════════════
# 项目管理 — 按国家 → 项目分组展示 (卡片式 UI)
# ════════════════════════════════════════════════════════════════════════════

_PM_CSS = """
<style>
.pm-country-banner {
    background: linear-gradient(135deg, #1F4E79 0%, #2E75B6 100%);
    color: white; padding: 14px 24px; border-radius: 10px;
    margin: 20px 0 14px 0;
}
.pm-country-banner .pm-title { font-size: 1.35rem; font-weight: 700; margin: 0; }
.pm-country-banner .pm-sub { opacity: 0.85; font-size: 0.88rem; margin-top: 2px; }
.pm-group-bar {
    background: #f0f4f8; border-left: 4px solid #2E75B6;
    padding: 10px 18px; border-radius: 0 8px 8px 0;
    margin: 14px 0 10px 0; display: flex; align-items: center; gap: 12px;
}
.pm-group-bar .pm-gname { font-weight: 700; font-size: 1.05rem; color: #1F4E79; }
.pm-tag {
    display: inline-block; padding: 2px 10px; border-radius: 10px;
    font-size: 0.75rem; font-weight: 600; color: #fff;
}
.pm-tag-onshore { background: #548235; }
.pm-tag-offshore { background: #1F4E79; }
.pm-tag-count { background: #aab; color: #fff; }
.pm-summary {
    background: #fafbfc; border: 1px solid #e8ecf0; border-radius: 8px;
    padding: 8px 16px; margin: 4px 0 14px 0; display: flex;
    gap: 24px; flex-wrap: wrap; font-size: 0.85rem; color: #555;
}
.pm-summary b { color: #1F4E79; }
.pm-best-tag {
    display: inline-block; font-size: 0.68rem; font-weight: 800;
    padding: 1px 8px; border-radius: 8px; margin-left: 6px; vertical-align: middle;
}
.pm-best-irr { background: linear-gradient(135deg, #FFD700, #FFA500); color: #333; }
.pm-best-lcoe { background: linear-gradient(135deg, #90EE90, #2E8B57); color: #fff; }
</style>
"""


def _irr_color(v: float) -> str:
    if v >= 0.10:
        return "#548235"
    elif v >= 0.06:
        return "#BF8F00"
    return "#C00000"


def _diagnose_project(inp: WindFarmFinancialInputs, res: CalculationResult) -> list[tuple[str, str, str]]:
    """Return list of (icon, message, severity) for result health check.
    Thresholds are country-specific, based on local policy/benchmark data."""
    diags: list[tuple[str, str, str]] = []
    country = inp.basic.country
    ptype = inp.basic.project_type
    th = get_diag_thresholds(country)

    irr = res.project_irr_after_tax
    eq = res.equity_irr
    lcoe = res.lcoe
    tariff = inp.tax_financial.tariff_with_tax
    npv = res.project_npv_after_tax
    payback = res.payback_after_tax
    capex = inp.investment.resolve_unit_investment()
    p90 = inp.basic.full_load_hours
    disc = inp.tax_financial.discount_rate
    roe = res.roe if hasattr(res, "roe") else None

    _cn = th.country_cn

    # ── IRR (project after-tax) ──
    if irr < 0:
        diags.append(("🔴", f"全投IRR为负({irr:.2%})，项目亏损", "error"))
    elif irr < th.irr_floor:
        diags.append(("🔴", f"全投IRR仅{irr:.2%}，远低于{_cn}基准({th.irr_low:.0%})", "error"))
    elif irr < th.irr_low:
        _note = f" ({th.irr_note})" if th.irr_note else ""
        diags.append(("🟠", f"全投IRR {irr:.2%}，低于{_cn}行业基准{th.irr_low:.0%}{_note}", "warning"))
    elif irr > th.irr_high:
        diags.append(("🟡", f"全投IRR {irr:.2%}，高于{_cn}正常区间({th.irr_high:.0%}) — 核查CAPEX/P90", "warning"))

    # ── Equity IRR ──
    if eq < 0:
        diags.append(("🔴", f"资本金IRR为负({eq:.2%})，融资结构亏损", "error"))
    elif eq < th.eq_irr_floor:
        diags.append(("🔴", f"资本金IRR {eq:.2%}，远低于{_cn}基准", "error"))
    elif eq < th.eq_irr_low:
        _note = f" ({th.eq_irr_note})" if th.eq_irr_note else ""
        diags.append(("🟠", f"资本金IRR {eq:.2%}，低于{_cn}基准{th.eq_irr_low:.0%}{_note}", "warning"))
    elif eq > th.eq_irr_high:
        _note = f" ({th.eq_irr_note})" if th.eq_irr_note else ""
        diags.append(("🟡", f"资本金IRR {eq:.2%}，高于{_cn}正常上限{th.eq_irr_high:.0%}{_note}", "warning"))

    if 0 < eq < irr:
        diags.append(("⚠️", "资本金IRR < 全投IRR — 贷款利率高于项目收益，杠杆效应为负", "info"))

    # ── LCOE ──
    lcoe_limit = th.lcoe_high_offshore if ptype == "offshore" else th.lcoe_high_onshore
    if lcoe > tariff > 0:
        diags.append(("🔴", f"LCOE({lcoe:.4f}) > 电价({tariff:.4f})，度电亏损", "error"))
    elif lcoe > lcoe_limit:
        diags.append(("🟠", f"LCOE {lcoe:.4f} 高于{_cn}{ptype}参考上限({lcoe_limit:.3f})", "warning"))

    # ── Tariff vs country range ──
    _tr = th.tariff_offshore if ptype == "offshore" else th.tariff_onshore
    if _tr[1] > 0:
        if tariff > _tr[1] * 1.3:
            diags.append(("🟡", f"电价{tariff:.4f}高于{_cn}政策区间({_tr[0]:.3f}-{_tr[1]:.3f}) — 确认电价来源", "warning"))
        elif tariff < _tr[0] * 0.7 and _tr[0] > 0:
            diags.append(("🟠", f"电价{tariff:.4f}低于{_cn}政策区间({_tr[0]:.3f}-{_tr[1]:.3f})", "warning"))

    # ── CAPEX ──
    capex_limit = th.capex_high_offshore if ptype == "offshore" else th.capex_high_onshore
    if capex > capex_limit * 1.3:
        diags.append(("🟡", f"CAPEX {capex:,.0f}$/kW 远超{_cn}{ptype}参考上限({capex_limit:,.0f}) — 核查EPC明细", "warning"))
    elif capex > capex_limit:
        diags.append(("🟠", f"CAPEX {capex:,.0f}$/kW 高于{_cn}{ptype}参考({capex_limit:,.0f})", "info"))

    # ── NPV ──
    if npv < 0:
        if irr > 0 and irr < disc:
            diags.append(("ℹ️", f"NPV {npv/1e6:,.1f}M$ (折现率{disc:.0%} > IRR{irr:.2%}，属正常)", "info"))
        else:
            diags.append(("🟠", f"NPV为负({npv/1e6:,.1f}M$)", "warning"))

    # ── Payback ──
    op_years = inp.operational.operation_years
    if payback > op_years:
        diags.append(("🔴", f"回收期{payback:.1f}年 > 运营期{op_years}年", "error"))
    elif payback > th.payback_critical:
        diags.append(("🟠", f"回收期{payback:.1f}年，偏长", "warning"))
    elif payback > th.payback_warning:
        diags.append(("ℹ️", f"回收期{payback:.1f}年", "info"))

    # ── P90 sanity ──
    if p90 > 4500:
        diags.append(("🟡", f"P90 {p90:.0f}h 异常高 — 核查风资源数据", "warning"))
    elif p90 < 1200 and ptype == "onshore":
        diags.append(("🟠", f"P90 {p90:.0f}h 偏低(陆上一般1800-3000h)", "warning"))

    return diags


def _render_project_card(pid: str, proj: dict, best_irr: bool, best_lcoe: bool):
    """Render one project card inside a st.container with border."""
    inp: WindFarmFinancialInputs = proj["inputs"]
    res: CalculationResult = proj["result"]
    irr_c = _irr_color(res.project_irr_after_tax)

    with st.container(border=True):
        # 标题行
        title_md = f"**{proj['name']}**"
        if best_irr:
            title_md += ' <span class="pm-best-tag pm-best-irr">Best IRR</span>'
        elif best_lcoe:
            title_md += ' <span class="pm-best-tag pm-best-lcoe">Best LCOE</span>'
        st.markdown(title_md, unsafe_allow_html=True)

        # 健康度诊断
        diags = _diagnose_project(inp, res)
        if diags:
            _severity_order = {"error": 0, "warning": 1, "info": 2}
            diags.sort(key=lambda d: _severity_order.get(d[2], 9))
            _diag_lines = " / ".join(f"{d[0]} {d[1]}" for d in diags[:3])
            _color = "#C00000" if diags[0][2] == "error" else "#BF8F00" if diags[0][2] == "warning" else "#666"
            st.markdown(f'<div style="font-size:0.78rem;color:{_color};margin:-8px 0 6px 0">{_diag_lines}</div>',
                        unsafe_allow_html=True)

        # KPI 指标
        k1, k2, k3, k4 = st.columns(4)
        k1.metric("Capacity", f"{inp.capacity_mw:.0f} MW")
        k2.metric("Investment", f"{inp.investment.resolve_unit_investment():,.0f} $/kW")
        k3.metric("P90", f"{inp.basic.full_load_hours:.0f} h")
        k4.metric("Tariff", f"{inp.tax_financial.tariff_with_tax:.4f} $/kWh")

        m1, m2, m3, m4 = st.columns(4)
        m1.metric("IRR (after tax)", f"{res.project_irr_after_tax:.2%}")
        m2.metric("Equity IRR", f"{res.equity_irr:.2%}")
        m3.metric("LCOE", f"{res.lcoe:.5f} $/kWh")
        m4.metric("NPV", f"{res.project_npv_after_tax / 1e6:,.1f} M$")

        r1, r2 = st.columns([1, 1])
        r1.metric("Payback", f"{res.payback_after_tax:.1f} yr")
        r2.metric("IRR (before tax)", f"{res.project_irr_before_tax:.2%}")

        # 操作按钮行
        b1, b2, b3, _pad = st.columns([1, 1, 1, 3])
        with b1:
            if st.button("View Details", key=f"view_{pid}", type="primary", use_container_width=True):
                st.session_state.detail_pid = pid
                st.rerun()
        with b2:
            excel_b = export_to_excel(inp, res)
            st.download_button(
                "Download Excel", data=excel_b,
                file_name=f"{proj['name']}_report.xlsx",
                mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                key=f"dl_{pid}", use_container_width=True,
            )
        with b3:
            confirming = st.session_state.get("confirm_delete") == pid
            if confirming:
                if st.button("✕ Cancel", key=f"cdel_n_{pid}", use_container_width=True):
                    st.session_state.pop("confirm_delete", None)
                    st.rerun()
            else:
                if st.button("🗑 Delete", key=f"del_{pid}", use_container_width=True):
                    st.session_state["confirm_delete"] = pid
                    st.rerun()

        # 删除确认区 — 全宽展示，避免挤压导致文字竖排
        if st.session_state.get("confirm_delete") == pid:
            st.warning(f"⚠️ 确认删除 **{proj['name']}**？此操作不可撤销，请输入管理员凭证：")
            _dc1, _dc2 = st.columns(2)
            del_u = _dc1.text_input("用户名", key=f"delu_{pid}", placeholder="Username")
            del_p = _dc2.text_input("密码", key=f"delp_{pid}", type="password", placeholder="Password")
            _db1, _db2, _pad2 = st.columns([1, 1, 4])
            with _db1:
                if st.button("✅ 确认删除", key=f"cdel_y_{pid}", type="primary", use_container_width=True):
                    if del_u == _DELETE_USER and del_p == _DELETE_PWD:
                        delete_project(pid)
                        st.rerun()
                    else:
                        st.error("用户名或密码错误")
            with _db2:
                if st.button("取消", key=f"cdel_n2_{pid}", use_container_width=True):
                    st.session_state.pop("confirm_delete", None)
                    st.rerun()


def _update_project_result(pid: str, proj: dict, new_inputs: WindFarmFinancialInputs):
    """Recalculate and update a project in-place + sync to DB."""
    new_result = calculate(new_inputs)
    proj["inputs"] = copy.deepcopy(new_inputs)
    proj["result"] = new_result
    if _USE_DB:
        try:
            _db.db_save(pid, proj["name"], proj.get("group", ""),
                        proj.get("country", ""), new_inputs,
                        proj.get("saved_at", time.strftime("%Y-%m-%d %H:%M:%S")))
        except Exception:
            pass


def _render_quick_edit(pid: str, proj: dict):
    """Quick parameter edit panel for a single project."""
    inp = proj["inputs"]
    with st.expander("🔧 快速修参 & 重新计算", expanded=False):
        st.caption("修改以下参数后点击「重新计算」，结果自动更新并同步到云端。")
        _c1, _c2, _c3, _c4 = st.columns(4)
        _new_p90 = _c1.number_input("P90 满负荷小时(h)", 500, 6000,
                                     int(inp.basic.full_load_hours), step=10, key=f"qe_p90_{pid}")
        _new_tariff = _c2.number_input("含税电价($/kWh)", 0.001, 0.50,
                                        float(inp.tax_financial.tariff_with_tax), step=0.001,
                                        format="%.4f", key=f"qe_tariff_{pid}")
        _new_capex = _c3.number_input("单位投资($/kW)", 100.0, 5000.0,
                                       float(inp.investment.resolve_unit_investment()), step=10.0,
                                       format="%.0f", key=f"qe_capex_{pid}")
        _new_disc = _c4.number_input("折现率(%)", 1.0, 20.0,
                                      float(inp.tax_financial.discount_rate * 100), step=0.5,
                                      format="%.1f", key=f"qe_disc_{pid}")

        _fc1, _fc2, _fc3, _fc4 = st.columns(4)
        _new_eq_ratio = _fc1.number_input("资本金比例(%)", 10.0, 100.0,
                                           float(inp.financing.equity_ratio * 100), step=5.0,
                                           format="%.0f", key=f"qe_eq_{pid}")
        _new_loan_rate = _fc2.number_input("贷款利率(%)", 0.5, 15.0,
                                            float(inp.financing.long_term_loan_rate * 100), step=0.1,
                                            format="%.2f", key=f"qe_loan_{pid}")
        _new_loss = _fc3.number_input("损耗率(%)", 0.0, 20.0,
                                       float(inp.basic.loss_rate * 100), step=0.5,
                                       format="%.1f", key=f"qe_loss_{pid}")
        _new_tax = _fc4.number_input("所得税率(%)", 0.0, 40.0,
                                      float(inp.tax_financial.income_tax_rate * 100), step=1.0,
                                      format="%.0f", key=f"qe_itax_{pid}")

        if st.button("🔄 重新计算此方案", key=f"qe_recalc_{pid}", type="primary"):
            new_inp = copy.deepcopy(inp)
            new_inp.basic.full_load_hours = _new_p90
            new_inp.tax_financial.tariff_with_tax = _new_tariff
            new_inp.investment.unit_static_investment = _new_capex
            new_inp.tax_financial.discount_rate = _new_disc / 100.0
            new_inp.financing.equity_ratio = _new_eq_ratio / 100.0
            new_inp.financing.long_term_loan_rate = _new_loan_rate / 100.0
            new_inp.basic.loss_rate = _new_loss / 100.0
            new_inp.tax_financial.income_tax_rate = _new_tax / 100.0
            if new_inp.investment.onshore_detail is not None:
                _ratio = _new_capex / max(inp.investment.resolve_unit_investment(), 1)
                new_inp.investment.onshore_detail.equipment_and_installation *= _ratio
                new_inp.investment.onshore_detail.civil_works *= _ratio
            if new_inp.investment.offshore_detail is not None:
                _ratio = _new_capex / max(inp.investment.resolve_unit_investment(), 1)
                new_inp.investment.offshore_detail.oem.turbine_price_per_kw *= _ratio
            _update_project_result(pid, proj, new_inp)
            st.success("已重新计算并保存！")
            st.rerun()


def _render_group_batch_edit(country: str, group_name: str, items: list):
    """Batch edit shared parameters for an entire project group."""
    _gkey = f"{country}__{group_name}"
    _batch_key = f"batch_edit_{_gkey}"

    if st.session_state.get(_batch_key):
        with st.container(border=True):
            st.markdown(f"**批量修参：{group_name}**（修改公共参数，所有 {len(items)} 个方案同时重算）")
            st.caption("仅修改下方参数，各方案的机型/台数/P90等独有参数保持不变。")

            _bc1, _bc2, _bc3, _bc4 = st.columns(4)
            ref_inp = items[0][1]["inputs"]
            _b_tariff = _bc1.number_input("含税电价($/kWh)", 0.001, 0.50,
                                           float(ref_inp.tax_financial.tariff_with_tax), step=0.001,
                                           format="%.4f", key=f"be_tariff_{_gkey}")
            _b_disc = _bc2.number_input("折现率(%)", 1.0, 20.0,
                                         float(ref_inp.tax_financial.discount_rate * 100), step=0.5,
                                         format="%.1f", key=f"be_disc_{_gkey}")
            _b_eq = _bc3.number_input("资本金比例(%)", 10.0, 100.0,
                                       float(ref_inp.financing.equity_ratio * 100), step=5.0,
                                       format="%.0f", key=f"be_eq_{_gkey}")
            _b_loan = _bc4.number_input("贷款利率(%)", 0.5, 15.0,
                                          float(ref_inp.financing.long_term_loan_rate * 100), step=0.1,
                                          format="%.2f", key=f"be_loan_{_gkey}")

            _bc5, _bc6, _bc7, _bc8 = st.columns(4)
            _b_loss = _bc5.number_input("损耗率(%)", 0.0, 20.0,
                                         float(ref_inp.basic.loss_rate * 100), step=0.5,
                                         format="%.1f", key=f"be_loss_{_gkey}")
            _b_itax = _bc6.number_input("所得税率(%)", 0.0, 40.0,
                                          float(ref_inp.tax_financial.income_tax_rate * 100), step=1.0,
                                          format="%.0f", key=f"be_itax_{_gkey}")
            _b_vat = _bc7.number_input("增值税率(%)", 0.0, 20.0,
                                        float(ref_inp.tax_financial.vat_rate * 100), step=1.0,
                                        format="%.0f", key=f"be_vat_{_gkey}")
            _b_opyrs = _bc8.number_input("运营期(年)", 15, 40,
                                          int(ref_inp.operational.operation_years), step=1,
                                          key=f"be_opyrs_{_gkey}")

            _bb1, _bb2, _bbpad = st.columns([1, 1, 4])
            with _bb1:
                if st.button("🔄 批量重新计算", key=f"be_go_{_gkey}", type="primary", use_container_width=True):
                    _count = 0
                    for _pid, _proj in items:
                        new_inp = copy.deepcopy(_proj["inputs"])
                        new_inp.tax_financial.tariff_with_tax = _b_tariff
                        new_inp.tax_financial.discount_rate = _b_disc / 100.0
                        new_inp.financing.equity_ratio = _b_eq / 100.0
                        new_inp.financing.long_term_loan_rate = _b_loan / 100.0
                        new_inp.basic.loss_rate = _b_loss / 100.0
                        new_inp.tax_financial.income_tax_rate = _b_itax / 100.0
                        new_inp.tax_financial.vat_rate = _b_vat / 100.0
                        new_inp.operational.operation_years = _b_opyrs
                        _update_project_result(_pid, _proj, new_inp)
                        _count += 1
                    st.session_state.pop(_batch_key, None)
                    st.success(f"已批量重算 {_count} 个方案！")
                    st.rerun()
            with _bb2:
                if st.button("取消", key=f"be_cancel_{_gkey}", use_container_width=True):
                    st.session_state.pop(_batch_key, None)
                    st.rerun()


def _render_project_list():
    """按 国家 -> 项目组 -> 方案 三级结构展示"""
    st.markdown(_PM_CSS, unsafe_allow_html=True)

    projects = st.session_state.projects

    from collections import defaultdict
    tree: dict[str, dict[str, list]] = defaultdict(lambda: defaultdict(list))
    for pid, proj in projects.items():
        country = proj.get("country", proj["inputs"].basic.country)
        group = proj.get("group", proj["name"].split(" - ")[0].strip() if " - " in proj["name"] else "Other")
        tree[country][group].append((pid, proj))

    total_variants = len(projects)
    total_groups = sum(len(g) for g in tree.values())

    # ── 筛选栏 ──
    _fc1, _fc2, _fc3 = st.columns([2, 4, 2])
    with _fc1:
        all_countries = ["全部"] + sorted(tree.keys())
        _sel_country = st.selectbox(
            "筛选国家", all_countries, index=0, key="pm_country_filter")
    with _fc2:
        _search = st.text_input(
            "搜索项目名称", value="", placeholder="输入关键词快速定位...",
            key="pm_search")
    with _fc3:
        st.markdown(
            f'<div class="pm-summary" style="margin-top:24px">'
            f'<span><b>{total_variants}</b> variants</span>'
            f'<span><b>{len(tree)}</b> countries</span></div>',
            unsafe_allow_html=True,
        )

    filtered_countries = sorted(tree.keys())
    if _sel_country != "全部":
        filtered_countries = [c for c in filtered_countries if c == _sel_country]

    for country in filtered_countries:
        groups = tree[country]
        country_count = sum(len(v) for v in groups.values())
        st.markdown(
            f'<div class="pm-country-banner">'
            f'<div class="pm-title">{country}</div>'
            f'<div class="pm-sub">{country_count} variants / {len(groups)} project groups</div>'
            f'</div>',
            unsafe_allow_html=True,
        )

        for group_name in sorted(groups.keys()):
            items = groups[group_name]
            if _search:
                items = [(pid, p) for pid, p in items
                         if _search.lower() in p["name"].lower()]
                if not items:
                    continue
            items.sort(key=lambda x: x[1]["name"])

            ptype = items[0][1]["inputs"].basic.project_type
            tag_cls = "pm-tag-onshore" if ptype == "onshore" else "pm-tag-offshore"
            type_label = "Onshore" if ptype == "onshore" else "Offshore"

            best_irr_pid = max(items, key=lambda x: x[1]["result"].project_irr_after_tax)[0]
            best_lcoe_pid = min(items, key=lambda x: x[1]["result"].lcoe)[0]

            _gkey = f"{country}__{group_name}"
            _rename_key = f"renaming_group_{_gkey}"
            _delgrp_key = f"delgroup_{_gkey}"

            # Group bar + 管理按钮
            _batch_key = f"batch_edit_{_gkey}"
            _gb_col, _ge_col, _gr_col, _gd_col = st.columns([6, 1, 1, 1])
            with _gb_col:
                st.markdown(
                    f'<div class="pm-group-bar" style="margin-top:8px">'
                    f'<span class="pm-gname">{group_name}</span>'
                    f'<span class="pm-tag {tag_cls}">{type_label}</span>'
                    f'<span class="pm-tag pm-tag-count">{len(items)} variants</span>'
                    f'</div>',
                    unsafe_allow_html=True,
                )
            with _ge_col:
                st.markdown("<div style='margin-top:16px'></div>", unsafe_allow_html=True)
                if st.button("🔄 批量修参", key=f"gbatch_{_gkey}", use_container_width=True):
                    st.session_state[_batch_key] = True
                    st.session_state.pop(_rename_key, None)
                    st.session_state.pop(_delgrp_key, None)
                    st.rerun()
            with _gr_col:
                st.markdown("<div style='margin-top:16px'></div>", unsafe_allow_html=True)
                if st.button("✏️ 重命名", key=f"grename_{_gkey}", use_container_width=True):
                    st.session_state[_rename_key] = True
                    st.session_state.pop(_delgrp_key, None)
                    st.session_state.pop(_batch_key, None)
                    st.rerun()
            with _gd_col:
                st.markdown("<div style='margin-top:16px'></div>", unsafe_allow_html=True)
                if st.button("🗑 删除组", key=f"gdelete_{_gkey}", use_container_width=True):
                    st.session_state[_delgrp_key] = True
                    st.session_state.pop(_rename_key, None)
                    st.session_state.pop(_batch_key, None)
                    st.rerun()

            # ── 批量修参 ──
            _render_group_batch_edit(country, group_name, items)

            # ── 重命名表单 ──
            if st.session_state.get(_rename_key):
                with st.container(border=True):
                    st.markdown(f"**重命名项目组：{group_name}**（所有方案的 group 和名称前缀同步更新）")
                    _rn1, _rn2 = st.columns(2)
                    _new_gname = _rn1.text_input("新组名", value=group_name, key=f"new_gname_{_gkey}")
                    _rn2.markdown("")
                    _ra1, _ra2 = st.columns(2)
                    _gadm_u = _ra1.text_input("管理员用户名", key=f"gadmu_{_gkey}", placeholder="Username")
                    _gadm_p = _ra2.text_input("管理员密码", key=f"gadmp_{_gkey}", type="password", placeholder="Password")
                    _rb1, _rb2, _rbpad = st.columns([1, 1, 4])
                    with _rb1:
                        if st.button("✅ 确认重命名", key=f"grn_ok_{_gkey}", type="primary", use_container_width=True):
                            if _gadm_u == _DELETE_USER and _gadm_p == _DELETE_PWD:
                                _new = _new_gname.strip()
                                if _new and _new != group_name:
                                    for _pid, _proj in items:
                                        _proj["group"] = _new
                                        if _proj["name"].startswith(f"{group_name} - "):
                                            _proj["name"] = _new + _proj["name"][len(group_name):]
                                            _proj["inputs"].basic.project_name = _proj["name"]
                                        if _USE_DB:
                                            try:
                                                _db.db_save(
                                                    _pid, _proj["name"], _new,
                                                    _proj.get("country", ""), _proj["inputs"],
                                                    _proj.get("saved_at", ""),
                                                )
                                            except Exception:
                                                pass
                                st.session_state.pop(_rename_key, None)
                                st.rerun()
                            else:
                                st.error("用户名或密码错误")
                    with _rb2:
                        if st.button("取消", key=f"grn_no_{_gkey}", use_container_width=True):
                            st.session_state.pop(_rename_key, None)
                            st.rerun()

            # ── 删除整组表单 ──
            if st.session_state.get(_delgrp_key):
                with st.container(border=True):
                    st.warning(f"⚠️ 将删除 **{group_name}** 组内全部 **{len(items)}** 个方案，操作不可撤销！")
                    _da1, _da2 = st.columns(2)
                    _dadm_u = _da1.text_input("管理员用户名", key=f"dadmu_{_gkey}", placeholder="Username")
                    _dadm_p = _da2.text_input("管理员密码", key=f"dadmp_{_gkey}", type="password", placeholder="Password")
                    _dc1, _dc2, _dcpad = st.columns([1, 1, 4])
                    with _dc1:
                        if st.button("✅ 确认全部删除", key=f"gdel_ok_{_gkey}", type="primary", use_container_width=True):
                            if _dadm_u == _DELETE_USER and _dadm_p == _DELETE_PWD:
                                for _pid, _ in items:
                                    delete_project(_pid)
                                st.session_state.pop(_delgrp_key, None)
                                st.rerun()
                            else:
                                st.error("用户名或密码错误")
                    with _dc2:
                        if st.button("取消", key=f"gdel_no_{_gkey}", use_container_width=True):
                            st.session_state.pop(_delgrp_key, None)
                            st.rerun()

            # 汇总对比表
            rows = []
            for pid, proj in items:
                inp = proj["inputs"]
                res = proj["result"]
                rows.append({
                    "Variant": proj["name"],
                    "MW": f"{inp.capacity_mw:.0f}",
                    "USD/kW": f"{inp.investment.resolve_unit_investment():,.0f}",
                    "P90 (h)": f"{inp.basic.full_load_hours:.0f}",
                    "Tariff": f"{inp.tax_financial.tariff_with_tax:.4f}",
                    "IRR pre-tax": f"{res.project_irr_before_tax:.2%}",
                    "IRR post-tax": f"{res.project_irr_after_tax:.2%}",
                    "Equity IRR": f"{res.equity_irr:.2%}",
                    "LCOE": f"{res.lcoe:.5f}",
                    "NPV (M$)": f"{res.project_npv_after_tax / 1e6:,.1f}",
                    "Payback (yr)": f"{res.payback_after_tax:.1f}",
                })
            st.dataframe(pd.DataFrame(rows), use_container_width=True, hide_index=True)

            # 对比图表 (方案 >= 2)
            if len(items) >= 2:
                names = [p[1]["name"].replace("Laguna-", "").replace("Laguna ", "") for p in items]
                irrs = [p[1]["result"].project_irr_after_tax * 100 for p in items]
                lcoes = [p[1]["result"].lcoe for p in items]
                invests = [p[1]["inputs"].investment.resolve_unit_investment() for p in items]

                col_c1, col_c2 = st.columns(2)
                with col_c1:
                    fig = go.Figure()
                    colors = [("#548235" if v >= 10 else "#BF8F00" if v >= 6 else "#C00000") for v in irrs]
                    fig.add_trace(go.Bar(
                        x=names, y=irrs, marker_color=colors,
                        text=[f"{v:.1f}%" for v in irrs], textposition="outside",
                    ))
                    fig.update_layout(
                        title="IRR post-tax (%)", height=300,
                        margin=dict(t=40, b=60, l=40, r=20),
                        yaxis_title="%", xaxis_tickangle=-35, font=dict(size=11),
                    )
                    st.plotly_chart(fig, use_container_width=True)
                with col_c2:
                    fig = go.Figure()
                    fig.add_trace(go.Bar(
                        x=names, y=invests, marker_color="#2E75B6", name="Investment",
                        text=[f"{v:,.0f}" for v in invests], textposition="outside",
                    ))
                    fig.add_trace(go.Scatter(
                        x=names, y=[l * 1e6 for l in lcoes], yaxis="y2",
                        name="LCOE", line=dict(color="#BF8F00", width=2.5),
                        marker=dict(size=8), mode="lines+markers+text",
                        text=[f"{l:.4f}" for l in lcoes], textposition="top center",
                    ))
                    fig.update_layout(
                        title="Investment & LCOE", height=300,
                        margin=dict(t=40, b=60, l=40, r=60),
                        yaxis=dict(title="USD/kW"),
                        yaxis2=dict(title="LCOE", overlaying="y", side="right", showgrid=False),
                        xaxis_tickangle=-35, font=dict(size=11), showlegend=False,
                    )
                    st.plotly_chart(fig, use_container_width=True)

            # 每个方案一张卡片 (2列布局)
            COLS_PER_ROW = 2
            for row_start in range(0, len(items), COLS_PER_ROW):
                row_items = items[row_start:row_start + COLS_PER_ROW]
                cols = st.columns(COLS_PER_ROW)
                for col_idx, (pid, proj) in enumerate(row_items):
                    with cols[col_idx]:
                        _render_project_card(
                            pid, proj,
                            best_irr=(pid == best_irr_pid and len(items) > 1),
                            best_lcoe=(pid == best_lcoe_pid and pid != best_irr_pid and len(items) > 1),
                        )

            st.markdown("")


# ════════════════════════════════════════════════════════════════════════════
# 主入口
# ════════════════════════════════════════════════════════════════════════════

def main():
    # ── 顶部标题 + 明阳 logo ──
    _logo_dir = os.path.dirname(os.path.abspath(__file__))
    _my_logo = os.path.join(_logo_dir, "mingyang_logo.png")

    # 明阳 logo — 侧边栏顶部，留白避开顶栏
    st.sidebar.markdown("<div style='margin-top:2.5rem'></div>", unsafe_allow_html=True)
    if os.path.exists(_my_logo):
        st.sidebar.image(_my_logo, width=220)
        st.sidebar.caption("Powered by **MINGYANG**")

    st.title("🌬️ 风电项目经济性评估")
    st.caption("Wind Farm Financial Assessment Dashboard | 多项目管理 & 对比 | 货币: USD")

    # ── 侧边栏底部：作者信息（内嵌 SVG 矢量猫 + 文字） ──
    st.sidebar.markdown("---")
    _cat_svg_inline = (
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 200 200" width="56" height="56">'
        '<polygon points="40,85 60,20 85,75" fill="#1a1a1a"/>'
        '<polygon points="160,85 140,20 115,75" fill="#1a1a1a"/>'
        '<polygon points="50,80 65,32 82,72" fill="#2a2a2a"/>'
        '<polygon points="150,80 135,32 118,72" fill="#2a2a2a"/>'
        '<ellipse cx="100" cy="105" rx="65" ry="55" fill="#1a1a1a"/>'
        '<ellipse cx="72" cy="100" rx="20" ry="22" fill="#fff"/>'
        '<ellipse cx="128" cy="100" rx="20" ry="22" fill="#fff"/>'
        '<ellipse cx="75" cy="102" rx="14" ry="16" fill="#111"/>'
        '<ellipse cx="125" cy="102" rx="14" ry="16" fill="#111"/>'
        '<circle cx="80" cy="95" r="5" fill="#fff"/>'
        '<circle cx="130" cy="95" r="5" fill="#fff"/>'
        '<circle cx="72" cy="105" r="2.5" fill="#fff"/>'
        '<circle cx="122" cy="105" r="2.5" fill="#fff"/>'
        '<ellipse cx="100" cy="120" rx="4" ry="3" fill="#333"/>'
        '<path d="M92,125 Q96,132 100,125" fill="none" stroke="#333" stroke-width="1.5" stroke-linecap="round"/>'
        '<path d="M100,125 Q104,132 108,125" fill="none" stroke="#333" stroke-width="1.5" stroke-linecap="round"/>'
        '<line x1="30" y1="108" x2="62" y2="115" stroke="#444" stroke-width="1.2"/>'
        '<line x1="28" y1="118" x2="60" y2="120" stroke="#444" stroke-width="1.2"/>'
        '<line x1="32" y1="128" x2="63" y2="124" stroke="#444" stroke-width="1.2"/>'
        '<line x1="170" y1="108" x2="138" y2="115" stroke="#444" stroke-width="1.2"/>'
        '<line x1="172" y1="118" x2="140" y2="120" stroke="#444" stroke-width="1.2"/>'
        '<line x1="168" y1="128" x2="137" y2="124" stroke="#444" stroke-width="1.2"/>'
        '<ellipse cx="100" cy="170" rx="45" ry="30" fill="#1a1a1a"/>'
        '<path d="M145,168 Q170,155 165,135 Q162,125 155,130" fill="none" stroke="#1a1a1a" stroke-width="10" stroke-linecap="round"/>'
        '</svg>'
    )
    _db_icon = "🟢" if _USE_DB else "⚪"
    _db_label = "Cloud DB" if _USE_DB else "Local Only"
    st.sidebar.markdown(
        "<div style='display:flex;align-items:center;gap:10px;padding:4px 0'>"
        f"<div style='width:56px;min-width:56px'>{_cat_svg_inline}</div>"
        "<div style='color:#888;line-height:1.4;font-size:0.9rem'>"
        "<b>MingYang Wind Tool</b><br>"
        f"Built by kurochilli &nbsp;{_db_icon} {_db_label}</div>"
        "</div>",
        unsafe_allow_html=True,
    )

    page = st.tabs(["📈 项目评估", "🔄 反算工具", "📊 项目对比", "🗂️ 项目管理", "🌍 各国市场概览"])

    # ═══════════════════════════════════════════════════════
    # Tab 1: 项目评估（含侧边栏全分项编辑）
    # ═══════════════════════════════════════════════════════
    with page[0]:
        # 模式切换
        input_mode = st.sidebar.radio(
            "Input Mode", ["Quick (10 params)", "Detailed (full)", "📤 Smart Upload"],
            horizontal=True, key="input_mode",
            help="Quick: 10 params. Detailed: full edit. Smart Upload: Excel/image auto-extract.",
        )
        st.sidebar.markdown("---")

        if input_mode.startswith("📤"):
            smart_upload_panel()
            inputs = None
        elif input_mode.startswith("Quick"):
            inputs = sidebar_inputs_quick()
        else:
            inputs = sidebar_inputs()

        if inputs is not None:
            result = calculate(inputs)

            col_save, col_info = st.columns([1, 3])
            with col_save:
                if st.button("💾 保存当前项目", type="primary"):
                    pid = save_project(inputs.basic.project_name, inputs, result)
                    st.success(f"项目已保存! (ID: {pid})")
            with col_info:
                st.caption(f"当前已保存 {len(st.session_state.projects)} 个项目")

            st.markdown("---")
            render_full_assessment(inputs, result, key_prefix="main")

            st.session_state["_reverse_inputs"] = inputs

    # ═══════════════════════════════════════════════════════
    # Tab 2: 反算工具（独立面板）
    # ═══════════════════════════════════════════════════════
    with page[1]:
        st.header("🔄 反算工具 / Reverse Solver")
        st.markdown(
            "**使用方法**: 先在左侧栏填好项目参数（Quick / Detailed 模式），"
            "或在 [项目评估] Tab 完成一次计算。本模块自动读取左侧栏参数，"
            "您**只需输入目标值**，系统迭代求解未知量。"
        )
        if "_reverse_inputs" in st.session_state:
            reverse_calc_panel(st.session_state["_reverse_inputs"])
        elif "input_mode" in st.session_state and not st.session_state.get("input_mode", "").startswith("📤"):
            try:
                if st.session_state.get("input_mode", "").startswith("Quick"):
                    rev_inputs = sidebar_inputs_quick()
                else:
                    rev_inputs = sidebar_inputs()
                if rev_inputs is not None:
                    reverse_calc_panel(rev_inputs)
            except Exception:
                st.warning("Please complete project parameters in the sidebar first, then switch to this tab.")
        else:
            st.warning(
                "Please fill in project parameters in the sidebar (Quick or Detailed mode) "
                "and run a calculation in the [Project Assessment] tab first."
            )

    # ═══════════════════════════════════════════════════════
    # Tab 3: 项目对比
    # ═══════════════════════════════════════════════════════
    with page[2]:
        comparison_page()

    # ═══════════════════════════════════════════════════════
    # Tab 4: 项目管理
    # ═══════════════════════════════════════════════════════
    with page[3]:
        st.header("🗂️ 已保存的项目")

        if not st.session_state.projects:
            st.info("暂无已保存的项目。请在「项目评估」页面编辑参数后点击「保存当前项目」。")
        elif "detail_pid" in st.session_state and st.session_state.detail_pid in st.session_state.projects:
            dpid = st.session_state.detail_pid
            dproj = st.session_state.projects[dpid]
            dinp = dproj["inputs"]
            dres = dproj["result"]

            if st.button("⬅️ 返回项目列表", key="back_to_list"):
                del st.session_state.detail_pid
                st.rerun()

            st.markdown(f"## 📄 {dproj['name']}")
            st.caption(f"保存于 {dproj['saved_at']}")

            # ── 健康诊断 ──
            _detail_diags = _diagnose_project(dinp, dres)
            if _detail_diags:
                for _icon, _msg, _sev in _detail_diags:
                    if _sev == "error":
                        st.error(f"{_icon} {_msg}")
                    elif _sev == "warning":
                        st.warning(f"{_icon} {_msg}")
                    else:
                        st.info(f"{_icon} {_msg}")

            # ── 快速修参 & 重新计算 ──
            _render_quick_edit(dpid, dproj)

            render_full_assessment(dinp, dres, key_prefix=f"detail_{dpid}")
        else:
            _render_project_list()

    # ═══════════════════════════════════════════════════════
    # Tab 5: 各国市场概览
    # ═══════════════════════════════════════════════════════
    with page[4]:
        render_market_overview()


if __name__ == "__main__":
    main()
